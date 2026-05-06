"""Client Report tab — branded PPTX advertising report for renewals.

Renders the five-section advertising report (Benchmarking Scatter,
Benchmarking Summary, Job Postings, Advertising ROI, Media Performance)
with optional PowerPoint export driven by the Renewals.pptx template.
"""

import io
import re
from datetime import datetime, timedelta

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots

from pptx import Presentation

from theme.colors import JGP_COLORS, JGP_PLOTLY_TEMPLATE
from data.processing import apply_media_categories
from data.loader import load_client_hq_regions


# Static explainers describing how each chart is calculated. Single source of
# truth: rendered on-screen as captions and substituted into the PPTX template
# via {{chart_explainer_<key>}} placeholders. Wording change → one-line edit
# here, no PowerPoint round trip.
CHART_EXPLAINERS = {
    'benchmark_scatter': (
        "Each marker is one of your vacancies. Its position shows how views "
        "and applies compare to the average for the same occupation across "
        "all other clients — top-right is above benchmark on both."
    ),
    'benchmark_average': (
        "Your average views and applies per vacancy as a percentage of the "
        "wider market average. 100% sits in line with the benchmark; above "
        "100% outperforms it."
    ),
    'postings_by_type': (
        "How your vacancy volume and apply clicks distribute across "
        "occupation categories during the report period."
    ),
    'spend_vs_ratecard': (
        "Your subscription spend (purple) stacked with the saving versus "
        "paying rate-card per vacancy (green). The full bar is what these "
        "postings would cost without your subscription."
    ),
    'cost_per_app_by_occupation': (
        "Your annual spend allocated to each occupation by share of "
        "vacancies, then divided by the applies generated. Lower bars "
        "indicate where spend produces candidates most efficiently."
    ),
    'media_performance': (
        "Average views and applies per vacancy, broken down by traffic "
        "source. Shows which channels (organic search, paid, direct, "
        "referral, etc.) drive the most candidates."
    ),
    'salary_by_occupation': (
        "Market salary spread for your top 10 most-priced occupations. "
        "Lines mark your mean (red), the national mean (amber) and your "
        "HQ-region mean (green). Ranked by count of priced vacancies; "
        "minimum 5 per occupation."
    ),
}


def generate_section_commentary(section, data):
    """Generate template-based commentary for each report section.

    Args:
        section: One of 'scatter', 'benchmark', 'postings', 'roi', 'media'.
        data: Dict with relevant metrics for the section.

    Returns:
        Markdown-formatted string with 2-4 sentences of data-driven insight.
    """
    fallback = "Insufficient data for detailed commentary."

    if section == 'scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        worst_performers = data.get('worst_performers', [])

        if total == 0:
            return fallback
        if total == 1:
            return ("Only one vacancy is available for analysis, which limits "
                    "benchmarking comparisons. A larger sample of roles will "
                    "enable more meaningful performance insights.")

        parts = []
        parts.append(
            f"Of **{total}** vacancies, **{benchmarkable}** "
            f"({'all' if benchmarkable == total else f'{benchmarkable / total:.0%}'}) "
            f"have sufficient market data for benchmarking"
            f"{f', while **{no_benchmark}** lack benchmark data due to low sample sizes in their occupation category' if no_benchmark > 0 else ''}."
        )

        if top_performers:
            top = top_performers[0]
            parts.append(
                f"**{top['title']}** ({top['occupation']}) is a standout performer, "
                f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views "
                f"and {top['applies_diff_pct']:+.0f}% on applies."
            )

        if zero_applies > 0:
            pct_zero = zero_applies / total * 100
            parts.append(
                f"**{zero_applies}** role{'s' if zero_applies != 1 else ''} "
                f"({pct_zero:.0f}% of total) received zero apply clicks and "
                f"may benefit from revised job descriptions or enhanced visibility."
            )

        return " ".join(parts)

    elif section == 'benchmark':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return fallback

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        views_word = "more" if views_diff >= 0 else "fewer"

        parts = []
        parts.append(
            f"Your vacancies received **{abs(views_diff):.0f}% {views_word} views** "
            f"than the market average, "
            f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}."
        )

        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            parts.append(
                f"Apply rates sit **{abs(applies_diff):.0f}% {applies_word} benchmark**"
                f"{', reflecting strong candidate engagement' if applies_diff >= 10 else ', suggesting job descriptions may benefit from enhancement' if applies_diff < -10 else ', broadly in line with market expectations'}."
            )

        parts.append(
            f"This analysis is based on **{num_jobs}** vacancies posted by {client_name} "
            f"during the report period."
        )
        return " ".join(parts)

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0:
            return fallback
        if by_type is None or len(by_type) == 0:
            return f"{client_name} posted **{total_jobs}** vacancies with **{total_applies:,}** total apply clicks during this period."

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        parts = []

        if len(by_type) == 1:
            parts.append(
                f"All **{total_jobs}** vacancies fall under **{top['occupation']}**, "
                f"generating **{int(top['apply_clicks']):,}** apply clicks."
            )
        else:
            parts.append(
                f"**{top['occupation']}** leads apply generation with "
                f"**{int(top['apply_clicks']):,}** apply clicks across "
                f"**{int(top['jobs_posted'])}** postings."
            )
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                parts.append(
                    f"**{bottom['occupation']}** received no apply clicks despite "
                    f"**{int(bottom['jobs_posted'])}** postings — these roles may "
                    f"benefit from revised job titles or enhanced descriptions."
                )
            elif len(by_type) > 1:
                parts.append(
                    f"**{bottom['occupation']}** generated the fewest applies "
                    f"(**{int(bottom['apply_clicks']):,}**) and may warrant "
                    f"targeted improvements to boost candidate engagement."
                )

        if total_applies == 0:
            parts.append(
                "No apply clicks were recorded across any category — reviewing "
                "listing quality and distribution channels is recommended."
            )
        return " ".join(parts)

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return "Enter your annual spend and rate card price to generate ROI commentary."

        parts = []
        if saving_pct > 0:
            parts.append(
                f"Your advertising investment delivers a **{saving_pct:.0f}% saving** "
                f"compared to rate card pricing, demonstrating strong value from "
                f"the platform partnership."
            )
        else:
            parts.append(
                "Current spend exceeds rate card value — reviewing the pricing "
                "structure or consolidating lower-performing listings may improve "
                "overall return on investment."
            )

        total_applies = data.get('total_applies', 0)
        if total_applies > 0:
            parts.append(
                f"At **£{cost_per_apply:,.2f} per apply**, each candidate "
                f"enquiry represents a cost-effective acquisition channel."
            )

        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            parts.append(
                f"**{best['occupation']}** achieves the best cost efficiency at "
                f"£{best['cost_per_apply']:,.2f} per apply, while "
                f"**{worst['occupation']}** is the most expensive at "
                f"£{worst['cost_per_apply']:,.2f}."
            )

        return " ".join(parts)

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return "Media source data is not yet available for this client."

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        parts = []

        parts.append(
            f"**{top_source['source_category']}** is the leading traffic source, "
            f"generating **{int(top_source['total_applies']):,}** apply clicks "
            f"from **{int(top_source['total_clicks']):,}** views."
        )

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        if best_conv['conversion_rate'] > 0:
            parts.append(
                f"**{best_conv['source_category']}** achieves the highest "
                f"view-to-apply conversion rate at **{best_conv['conversion_rate']:.1f}%**."
            )

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            parts.append(
                f"Paid channels ({paid['source_category']}) contributed "
                f"**{int(paid['total_applies']):,}** applies with a "
                f"**{paid['conversion_rate']:.1f}%** conversion rate."
            )

        return " ".join(parts)

    return fallback


def generate_section_commentary_structured(section, data):
    """Structured commentary for PPTX template — returns dict with intro + bullet points.

    Returns: dict with keys 'intro', 'point_1', 'point_2', 'point_3' (last is optional).
    Each value is plain text (no markdown). Empty string for unused points.
    """
    def _clean(text):
        """Strip markdown bold markers."""
        return re.sub(r'\*\*(.+?)\*\*', r'\1', text or '')

    if section == 'benchmark_scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        client_name = data.get('client_name', 'This client')

        if total == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': '', 'point_3': ''}

        intro = (f"{client_name} posted {total} vacancies during this period. "
                 f"Of these, {benchmarkable} have sufficient market data to benchmark against comparable public sector roles.")

        point_1 = ''
        if top_performers:
            top = top_performers[0]
            point_1 = (f"{top['title']} ({top['occupation']}) is the standout performer, "
                       f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views and "
                       f"{top['applies_diff_pct']:+.0f}% on applies.")

        point_2 = ''
        if zero_applies > 0:
            pct = zero_applies / total * 100 if total > 0 else 0
            point_2 = (f"{zero_applies} role{'s' if zero_applies != 1 else ''} ({pct:.0f}% of total) "
                       f"received zero apply clicks — these may benefit from revised job descriptions or enhanced visibility.")

        point_3 = ''
        if no_benchmark > 0:
            point_3 = (f"{no_benchmark} role{'s' if no_benchmark != 1 else ''} could not be benchmarked "
                       f"due to low market sample sizes in their occupation category.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    elif section == 'benchmark_average':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        intro = (f"Across {num_jobs} vacancies, {client_name} averaged {client_clicks:,.0f} views "
                 f"and {client_applies:,.1f} applies per role.")

        views_word = "more" if views_diff >= 0 else "fewer"
        point_1 = (f"Your vacancies received {abs(views_diff):.0f}% {views_word} views than the market average — "
                   f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}.")

        point_2 = ''
        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            tone = ('reflecting strong candidate engagement' if applies_diff >= 10
                    else 'suggesting job descriptions may benefit from enhancement' if applies_diff < -10
                    else 'broadly in line with market expectations')
            point_2 = f"Apply rates sit {abs(applies_diff):.0f}% {applies_word} benchmark, {tone}."

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0 or by_type is None or len(by_type) == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        intro = (f"{client_name} posted {total_jobs} vacancies generating {total_applies:,} apply clicks across "
                 f"{len(by_type)} occupation categories.")

        point_1 = (f"{top['occupation']} leads apply generation with {int(top['apply_clicks']):,} apply clicks across "
                   f"{int(top['jobs_posted'])} postings — your strongest performing category.")

        point_2 = ''
        if len(by_type) > 1:
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                point_2 = (f"{bottom['occupation']} received no apply clicks despite {int(bottom['jobs_posted'])} postings — "
                           f"these roles may benefit from revised titles or enhanced descriptions.")
            else:
                point_2 = (f"{bottom['occupation']} generated the fewest applies ({int(bottom['apply_clicks']):,}) "
                           f"and may warrant targeted improvements.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        rate_card_total = data.get('rate_card_total', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return {'intro': 'Enter your annual spend and rate card price to generate ROI commentary.',
                    'point_1': '', 'point_2': ''}

        saving_amount = rate_card_total - annual_spend
        intro = (f"Across {num_jobs} vacancies, your subscription delivered £{saving_amount:,.0f} of value "
                 f"compared to rate card pricing — a {saving_pct:.0f}% saving.")

        point_1 = (f"At £{cost_per_apply:,.2f} per apply, each candidate enquiry represents a "
                   f"cost-effective acquisition channel for {data.get('client_name', 'your team')}.")

        point_2 = ''
        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            point_2 = (f"{best['occupation']} achieves the best cost efficiency at £{best['cost_per_apply']:,.2f} per apply, "
                       f"while {worst['occupation']} is the most expensive at £{worst['cost_per_apply']:,.2f}.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return {'intro': 'Media source data is not yet available for this client.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        intro = (f"{client_name}'s vacancies received traffic from {len(cat_stats)} distinct channels. "
                 f"{top_source['source_category']} is the leading source, generating "
                 f"{int(top_source['total_applies']):,} applies from {int(top_source['total_clicks']):,} views.")

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        point_1 = ''
        if best_conv['conversion_rate'] > 0:
            point_1 = (f"{best_conv['source_category']} achieves the highest view-to-apply conversion rate at "
                       f"{best_conv['conversion_rate']:.1f}%, indicating well-matched candidates from this channel.")

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        point_2 = ''
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            point_2 = (f"Paid channels ({paid['source_category']}) contributed {int(paid['total_applies']):,} applies "
                       f"with a {paid['conversion_rate']:.1f}% conversion rate.")

        point_3 = ''
        if len(cat_stats) > 1:
            second_source = sorted_stats.iloc[1] if len(sorted_stats) > 1 else None
            if second_source is not None:
                point_3 = (f"{second_source['source_category']} is the second strongest channel with "
                           f"{int(second_source['total_applies']):,} applies — providing diversified candidate flow.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    elif section == 'salary':
        per_occ = data.get('per_occ') or []
        client_name = data.get('client_name', 'This client')
        client_region = data.get('client_region')

        if not per_occ:
            return {'intro': 'Insufficient salary data to generate commentary for this client.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        # Compare each occupation's client mean to the national mean.
        # Tuples: (occupation, pct_diff_signed, client_mean, national_mean)
        deltas = []
        for p in per_occ:
            c, n = p.get('client_mean'), p.get('national_mean')
            if c is None or n is None or pd.isna(c) or pd.isna(n) or n == 0:
                continue
            deltas.append((p['occupation'], (c - n) / n * 100, c, n))

        above = sorted([d for d in deltas if d[1] > 0], key=lambda x: x[1], reverse=True)
        below = sorted([d for d in deltas if d[1] < 0], key=lambda x: x[1])  # most-negative first

        n_total = len(deltas)
        if n_total == 0:
            return {'intro': 'Salary data was insufficient to compute market comparisons.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        intro = (f"Across {client_name}'s top {n_total} most-posted occupations with salary data, "
                 f"{len(above)} pay above the national market average and {len(below)} pay below it.")

        point_1 = ''
        if above:
            top = above[0]
            point_1 = (f"{top[0]} is your strongest premium — sitting {top[1]:.0f}% above the national "
                       f"average (£{top[2]:,.0f} vs £{top[3]:,.0f}). Useful signal for attraction in this discipline.")

        point_2 = ''
        if below:
            worst = below[0]
            point_2 = (f"{worst[0]} sits {abs(worst[1]):.0f}% below the national average "
                       f"(£{worst[2]:,.0f} vs £{worst[3]:,.0f}) — a likely contributor to slower candidate flow in this category.")

        point_3 = ''
        if client_region:
            reg_above = reg_below = 0
            for p in per_occ:
                c, r = p.get('client_mean'), p.get('regional_mean')
                if c is None or r is None or pd.isna(c) or pd.isna(r):
                    continue
                if c > r:
                    reg_above += 1
                elif c < r:
                    reg_below += 1
            if reg_above + reg_below > 0:
                point_3 = (f"Within {client_region}, {client_name} pays above the regional average for "
                           f"{reg_above} of these roles and below for {reg_below} — useful context when "
                           f"benchmarking against employers competing for the same local talent pool.")
        else:
            point_3 = ("Regional benchmark unavailable for this client — common for central-government and "
                       "multi-site bodies whose vacancies span the UK.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    return {'intro': 'No commentary available.', 'point_1': '', 'point_2': '', 'point_3': ''}


def render_client_report(df, media_df=None):
    """Render the Client Report tab — branded PPTX advertising report for renewals."""
    st.header("Client Advertising Report")

    # --- Controls ---
    col_client, col_dates = st.columns(2)
    with col_client:
        # Client selector — uses organization_name (the actual client/employer)
        org_col = 'organization_name' if 'organization_name' in df.columns else 'importer_name'
        orgs = sorted(df[org_col].dropna().unique())
        # Filter out blanks and unknowns
        orgs = [o for o in orgs if o and str(o).strip() not in ('', 'Unknown', 'nan')]
        org_counts = df.groupby(org_col).size().to_dict()
        org_labels = [f"{name} ({org_counts.get(name, 0):,} vacancies)" for name in orgs]
        selected_idx = st.selectbox(
            "Select Client / Organisation", range(len(orgs)),
            format_func=lambda i: org_labels[i],
            key='report_client'
        )
        selected_client = orgs[selected_idx]
    with col_dates:
        min_date = df['first_event_date'].dropna().min()
        max_date = df['last_event_date'].dropna().max()
        if pd.notna(min_date) and pd.notna(max_date):
            min_d = min_date.date() if hasattr(min_date, 'date') else min_date
            max_d = max_date.date() if hasattr(max_date, 'date') else max_date
        else:
            min_d = datetime.now().date() - timedelta(days=365)
            max_d = datetime.now().date()
        report_dates = st.date_input("Report Period", [min_d, max_d], key='report_dates')

    with st.expander("Cost & Report Settings", expanded=False):
        cost_col1, cost_col2 = st.columns(2)
        with cost_col1:
            annual_spend = st.number_input(
                "Annual Spend (GBP)", value=0.0, step=100.0, format="%.2f",
                key='report_spend',
                help="Enter 0 to skip the ROI section"
            )
        with cost_col2:
            rate_card_price = st.number_input(
                "Rate Card Price per Job (GBP)", value=600.0, step=10.0, format="%.2f",
                key='report_rate_card'
            )
        settings_col1, settings_col2 = st.columns(2)
        with settings_col1:
            include_self = st.checkbox(
                "Include self in benchmark",
                value=False,
                key='report_include_self',
                help="When unchecked, the selected client is excluded from the benchmark average (recommended for fair comparison)"
            )
        with settings_col2:
            st.markdown("**Contact details** (for PDF)")
            contact_name = st.text_input("Account Manager", key='report_contact_name', placeholder="e.g. Jane Smith")
            contact_title = st.text_input("Title", key='report_contact_title', placeholder="e.g. Account Director")
            contact_email = st.text_input("Email", key='report_contact_email', placeholder="e.g. jane@jgp.co.uk")
            contact_phone = st.text_input("Phone", key='report_contact_phone', placeholder="e.g. 020 7946 0958")

    generate_clicked = st.button("Generate Report", type="primary", key='report_generate')

    if not generate_clicked and 'report_generated' not in st.session_state:
        st.info("Select a client and click **Generate Report** to build the advertising report.")
        return

    st.session_state['report_generated'] = True

    # --- Data preparation ---
    if len(report_dates) < 2:
        st.warning("Please select a start and end date.")
        return

    report_start, report_end = report_dates[0], report_dates[1]

    # Client data — filter on organization_name (or importer_name fallback)
    client_df = df[df[org_col] == selected_client].copy()
    client_df = client_df[
        (client_df['last_event_date'].dt.date >= report_start) &
        (client_df['first_event_date'].dt.date <= report_end)
    ]

    if len(client_df) == 0:
        st.warning(f"No vacancies found for **{selected_client}** in the selected date range.")
        return

    # Benchmark = market data in same date range (exclude self by default for fair comparison)
    bench_mask = (
        (df['last_event_date'].dt.date >= report_start) &
        (df['first_event_date'].dt.date <= report_end)
    )
    if not include_self:
        bench_mask = bench_mask & (df[org_col] != selected_client)
    benchmark_df = df[bench_mask].copy()

    # Media data for client — match via entity_id (most reliable link)
    # Note: prepare_enriched_data() renames entity_id_str → entity_id in df,
    # but media_df keeps the original entity_id_str column name.
    client_media = None
    if media_df is not None and len(media_df) > 0:
        client_eid_col = 'entity_id' if 'entity_id' in client_df.columns else 'entity_id_str'
        media_eid_col = 'entity_id_str' if 'entity_id_str' in media_df.columns else 'entity_id'
        if client_eid_col in client_df.columns and media_eid_col in media_df.columns:
            client_entity_ids = client_df[client_eid_col].dropna().unique()
            client_media = media_df[media_df[media_eid_col].isin(client_entity_ids)].copy()
        if client_media is None or len(client_media) == 0:
            client_media = None
        else:
            client_media = apply_media_categories(client_media)

    # Store all figures for PDF export
    report_figures = {}

    st.markdown("---")

    # ===================================================================
    # SECTION 1: BENCHMARKING SCATTER
    # ===================================================================
    st.subheader("Benchmarking Jobs")

    # Calculate per-occupation benchmark averages (from ALL clients)
    occ_benchmarks = benchmark_df.groupby('occupation').agg(
        avg_clicks=('clicks', 'mean'),
        avg_applies=('applies', 'mean'),
        vacancy_count=('clicks', 'count')
    ).reset_index()

    # Only use occupations with enough data for reliable benchmarks
    MIN_BENCHMARK_VACANCIES = 5
    reliable_occs = occ_benchmarks[occ_benchmarks['vacancy_count'] >= MIN_BENCHMARK_VACANCIES]

    # Calculate % difference from benchmark for each client vacancy (vectorized)
    scatter_df = client_df[['title', 'occupation', 'clicks', 'applies']].copy()
    scatter_df = scatter_df.merge(
        reliable_occs[['occupation', 'avg_clicks', 'avg_applies']],
        on='occupation', how='left'
    )

    # Vectorized % diff calculations (safe division)
    scatter_df['views_diff_pct'] = (
        (scatter_df['clicks'] - scatter_df['avg_clicks'])
        / scatter_df['avg_clicks'].replace(0, np.nan) * 100
    ).fillna(0)
    scatter_df['applies_diff_pct'] = (
        (scatter_df['applies'] - scatter_df['avg_applies'])
        / scatter_df['avg_applies'].replace(0, np.nan) * 100
    ).fillna(0)

    # Categorize vacancies into 4 groups (matching original PDF categories)
    has_benchmark = scatter_df['avg_clicks'].notna()
    zero_applies = scatter_df['applies'] == 0
    occ_median_clicks = scatter_df.loc[has_benchmark, 'clicks'].median() if has_benchmark.any() else 0
    high_traffic = scatter_df['clicks'] > occ_median_clicks

    scatter_df['category'] = np.select(
        [has_benchmark & ~zero_applies,
         has_benchmark & zero_applies & high_traffic,
         has_benchmark & zero_applies & ~high_traffic],
        ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)'],
        default='Low Sample (No Benchmark)'
    )

    # Clean up merge columns and fill NaN diffs for no-benchmark rows
    scatter_df.drop(columns=['avg_clicks', 'avg_applies'], inplace=True)
    scatter_df.loc[~has_benchmark, ['views_diff_pct', 'applies_diff_pct']] = 0

    # --- Brand colour + marker maps for 4 categories ---
    category_colors = {
        'Benchmarkable': JGP_COLORS['primary'],
        'Zero Applies (Possible Redirect)': JGP_COLORS['amber'],
        'Zero Applies (Low Traffic)': JGP_COLORS['negative'],
        'Low Sample (No Benchmark)': JGP_COLORS['light_purple'],
    }
    category_symbols = {
        'Benchmarkable': 'triangle-up',
        'Zero Applies (Possible Redirect)': 'diamond',
        'Zero Applies (Low Traffic)': 'x',
        'Low Sample (No Benchmark)': 'circle',
    }

    chart_col, commentary_col = st.columns([3, 2])

    with chart_col:
        if len(scatter_df) > 0:
            fig_scatter = go.Figure()
            for cat in ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)', 'Low Sample (No Benchmark)']:
                cat_data = scatter_df[scatter_df['category'] == cat]
                if len(cat_data) == 0:
                    continue
                fig_scatter.add_trace(go.Scatter(
                    x=cat_data['applies_diff_pct'],
                    y=cat_data['views_diff_pct'],
                    mode='markers',
                    name=cat,
                    marker=dict(
                        color=category_colors[cat],
                        symbol=category_symbols[cat],
                        size=10,
                        line=dict(width=1, color='white')
                    ),
                    customdata=cat_data[['title', 'occupation', 'clicks', 'applies']].values,
                    hovertemplate=(
                        "<b>%{customdata[0]}</b><br>"
                        "Occupation: %{customdata[1]}<br>"
                        "Views: %{customdata[2]}<br>"
                        "Applies: %{customdata[3]}<br>"
                        "Views vs Bench: %{y:+.0f}%<br>"
                        "Applies vs Bench: %{x:+.0f}%<extra></extra>"
                    ),
                ))
            fig_scatter.add_hline(y=0, line_dash="dash", line_color="grey", opacity=0.5)
            fig_scatter.add_vline(x=0, line_dash="dash", line_color="grey", opacity=0.5)
            fig_scatter.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_scatter.update_layout(
                height=500,
                showlegend=True,
                legend=dict(orientation='h', y=-0.15, font=dict(size=11)),
                xaxis_title="Applies Difference from Benchmark (%)",
                yaxis_title="Views Difference from Benchmark (%)",
                plot_bgcolor='rgba(0,0,0,0)',
                xaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
                yaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
            )
            st.plotly_chart(fig_scatter, use_container_width=True)
            st.caption(CHART_EXPLAINERS['benchmark_scatter'])
            report_figures['scatter'] = fig_scatter

    with commentary_col:
        benchmarkable_count = len(scatter_df[scatter_df['category'] == 'Benchmarkable'])
        zero_redirect = len(scatter_df[scatter_df['category'] == 'Zero Applies (Possible Redirect)'])
        zero_low = len(scatter_df[scatter_df['category'] == 'Zero Applies (Low Traffic)'])
        no_bench_count = len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)'])

        kpi_c1, kpi_c2 = st.columns(2)
        with kpi_c1:
            st.metric("Benchmarkable", benchmarkable_count)
            st.metric("Zero Applies (Redirect)", zero_redirect)
        with kpi_c2:
            st.metric("Zero Applies (Low Traffic)", zero_low)
            st.metric("No Benchmark", no_bench_count)

        # Build top/worst performers for commentary
        benchmarkable_rows = scatter_df[scatter_df['category'] == 'Benchmarkable'].copy()
        top_performers = []
        worst_performers = []
        if len(benchmarkable_rows) > 0:
            benchmarkable_rows['_score'] = benchmarkable_rows['views_diff_pct'] + benchmarkable_rows['applies_diff_pct']
            top_sorted = benchmarkable_rows.nlargest(3, '_score')
            top_performers = top_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')
            worst_sorted = benchmarkable_rows.nsmallest(3, '_score')
            worst_performers = worst_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')

        commentary = generate_section_commentary('scatter', {
            'total_count': len(scatter_df),
            'benchmarkable_count': benchmarkable_count,
            'zero_applies_count': zero_redirect + zero_low,
            'no_benchmark_count': no_bench_count,
            'top_performers': top_performers,
            'worst_performers': worst_performers,
        })
        st.markdown(commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 2: BENCHMARKING SUMMARY
    # ===================================================================
    st.subheader("Benchmarking Summary")

    benchmark_avg_clicks = benchmark_df['clicks'].mean() if len(benchmark_df) > 0 else 0
    benchmark_avg_applies = benchmark_df['applies'].mean() if len(benchmark_df) > 0 else 0
    client_avg_clicks = client_df['clicks'].mean() if len(client_df) > 0 else 0
    client_avg_applies = client_df['applies'].mean() if len(client_df) > 0 else 0

    kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
    with kpi_col1:
        st.metric("Benchmark Avg. Views", f"{benchmark_avg_clicks:,.0f}")
    with kpi_col2:
        views_delta = client_avg_clicks - benchmark_avg_clicks
        st.metric("Your Jobs - Avg. Views", f"{client_avg_clicks:,.0f}",
                  delta=f"{views_delta:+,.0f} vs benchmark")
    with kpi_col3:
        st.metric("Benchmark Avg. Applies", f"{benchmark_avg_applies:,.1f}")
    with kpi_col4:
        applies_delta = client_avg_applies - benchmark_avg_applies
        st.metric("Your Jobs - Avg. Applies", f"{client_avg_applies:,.1f}",
                  delta=f"{applies_delta:+,.1f} vs benchmark")

    # Bar charts (60%) + commentary (40%)
    bench_chart_col, bench_commentary_col = st.columns([3, 2])

    with bench_chart_col:
        views_pct = (client_avg_clicks / benchmark_avg_clicks * 100) if benchmark_avg_clicks > 0 else 0
        applies_pct = (client_avg_applies / benchmark_avg_applies * 100) if benchmark_avg_applies > 0 else 0

        fig_bench = go.Figure()
        fig_bench.add_trace(go.Bar(
            x=['Views', 'Applies'],
            y=[views_pct, applies_pct],
            marker_color=[
                JGP_COLORS['positive'] if views_pct >= 100 else JGP_COLORS['amber'],
                JGP_COLORS['positive'] if applies_pct >= 100 else JGP_COLORS['amber'],
            ],
            text=[f"{views_pct:.0f}%", f"{applies_pct:.0f}%"],
            textposition='outside',
            textfont=dict(size=14, color=JGP_COLORS['deep_blue']),
        ))
        fig_bench.add_hline(
            y=100, line_dash="dash", line_width=3, line_color=JGP_COLORS['accent'],
            annotation_text="Benchmark (100%)",
            annotation_position="top right",
            annotation_bgcolor="white",
            annotation_bordercolor=JGP_COLORS['deep_blue'],
            annotation_borderwidth=1,
            annotation_font=dict(color=JGP_COLORS['deep_blue'], size=14),
        )
        fig_bench.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
        fig_bench.update_layout(
            title="Your Performance vs Market Benchmark",
            yaxis_title="% of Benchmark",
            height=400, showlegend=False,
            yaxis_range=[0, max(views_pct, applies_pct, 100) * 1.25],
            plot_bgcolor='rgba(0,0,0,0)',
        )
        st.plotly_chart(fig_bench, use_container_width=True)
        st.caption(CHART_EXPLAINERS['benchmark_average'])
        report_figures['benchmark_combined'] = fig_bench

    with bench_commentary_col:
        bench_commentary = generate_section_commentary('benchmark', {
            'client_avg_clicks': client_avg_clicks,
            'benchmark_avg_clicks': benchmark_avg_clicks,
            'client_avg_applies': client_avg_applies,
            'benchmark_avg_applies': benchmark_avg_applies,
            'num_jobs': len(client_df),
            'client_name': selected_client,
        })
        st.markdown(bench_commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 3: JOB POSTINGS BY TYPE
    # ===================================================================
    st.subheader("Job Postings by Type")

    by_type = client_df.groupby('occupation').agg(
        jobs_posted=('clicks', 'count'),
        apply_clicks=('applies', 'sum')
    ).reset_index().sort_values('jobs_posted', ascending=True)

    by_type = by_type[by_type['jobs_posted'] >= 1]

    chart_col3, commentary_col3 = st.columns([3, 2])

    with chart_col3:
        fig_postings = go.Figure()
        fig_postings.add_trace(go.Bar(
            y=by_type['occupation'], x=by_type['jobs_posted'],
            name='Jobs Posted', orientation='h',
            marker_color=JGP_COLORS['primary'],
            text=by_type['jobs_posted'], textposition='outside'
        ))
        fig_postings.add_trace(go.Bar(
            y=by_type['occupation'], x=by_type['apply_clicks'],
            name='Apply Clicks', orientation='h',
            marker_color=JGP_COLORS['accent'],
            text=by_type['apply_clicks'].astype(int), textposition='outside',
            textfont=dict(color=JGP_COLORS['deep_blue']),
        ))
        fig_postings.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
        fig_postings.update_layout(
            barmode='group', height=max(400, len(by_type) * 40),
            legend=dict(orientation='h', y=-0.1),
            xaxis_title="Count", yaxis_title="",
            plot_bgcolor='rgba(0,0,0,0)',
            bargap=0.1, bargroupgap=0.0,
        )
        st.plotly_chart(fig_postings, use_container_width=True)
        st.caption(CHART_EXPLAINERS['postings_by_type'])
        report_figures['postings'] = fig_postings

    with commentary_col3:
        total_jobs = len(client_df)
        total_applies_val = int(client_df['applies'].sum())
        postings_commentary = generate_section_commentary('postings', {
            'total_jobs': total_jobs,
            'total_applies': total_applies_val,
            'by_type': by_type,
            'client_name': selected_client,
        })
        st.markdown(postings_commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 4: ADVERTISING ROI
    # ===================================================================
    st.subheader("Advertising ROI")

    num_jobs = len(client_df)
    total_clicks = int(client_df['clicks'].sum())
    total_applies_val = int(client_df['applies'].sum())

    roi_by_type = None
    if annual_spend > 0:
        cost_per_job = annual_spend / num_jobs if num_jobs > 0 else 0
        cost_per_view = annual_spend / total_clicks if total_clicks > 0 else 0
        cost_per_apply = annual_spend / total_applies_val if total_applies_val > 0 else 0

        roi_kpi1, roi_kpi2, roi_kpi3, roi_kpi4 = st.columns(4)
        with roi_kpi1:
            st.metric("Jobs Advertised", f"{num_jobs:,}")
        with roi_kpi2:
            st.metric("Cost per Job", f"£{cost_per_job:,.2f}")
        with roi_kpi3:
            st.metric("Cost per View", f"£{cost_per_view:,.2f}")
        with roi_kpi4:
            st.metric("Cost per Apply", f"£{cost_per_apply:,.2f}")

        roi_chart_col, roi_commentary_col = st.columns([3, 2])

        with roi_chart_col:
            rate_card_total = rate_card_price * num_jobs
            saving_pct = ((rate_card_total - annual_spend) / rate_card_total * 100) if rate_card_total > 0 else 0

            fig_roi = go.Figure()
            fig_roi.add_trace(go.Bar(
                x=['Your Spend'], y=[annual_spend],
                name='Your Spend', marker_color=JGP_COLORS['primary'],
                text=[f"£{annual_spend:,.0f}"], textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_roi.add_trace(go.Bar(
                x=['Rate Card Value'], y=[rate_card_total],
                name='Rate Card Value', marker_color=JGP_COLORS['deep_blue'],
                text=[f"£{rate_card_total:,.0f}"], textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_roi.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_roi.update_layout(
                title=f"Cost vs Rate Card (Saving: {saving_pct:.0f}%)",
                height=350, showlegend=True, yaxis_title="GBP",
                plot_bgcolor='rgba(0,0,0,0)',
            )
            st.plotly_chart(fig_roi, use_container_width=True)
            st.caption(CHART_EXPLAINERS['spend_vs_ratecard'])
            report_figures['roi_cost'] = fig_roi

            # Cost per apply by type chart
            roi_by_type = client_df.groupby('occupation').agg(
                total_applies=('applies', 'sum'),
                job_count=('clicks', 'count')
            ).reset_index()
            roi_by_type = roi_by_type[roi_by_type['total_applies'] > 0]
            roi_by_type['cost_allocated'] = annual_spend * (roi_by_type['job_count'] / roi_by_type['job_count'].sum())
            roi_by_type['cost_per_apply'] = roi_by_type['cost_allocated'] / roi_by_type['total_applies']
            roi_by_type = roi_by_type.sort_values('cost_per_apply', ascending=True)

            if len(roi_by_type) > 0:
                fig_cpa = go.Figure()
                fig_cpa.add_trace(go.Bar(
                    y=roi_by_type['occupation'], x=roi_by_type['cost_per_apply'],
                    orientation='h', marker_color=JGP_COLORS['amber'],
                    text=roi_by_type['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                    textposition='outside'
                ))
                fig_cpa.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
                fig_cpa.update_layout(
                    title="Cost per Apply by Job Type",
                    height=max(300, len(roi_by_type) * 35),
                    xaxis_title="Cost per Apply (GBP)", yaxis_title="",
                    plot_bgcolor='rgba(0,0,0,0)',
                )
                st.plotly_chart(fig_cpa, use_container_width=True)
                st.caption(CHART_EXPLAINERS['cost_per_app_by_occupation'])
                report_figures['roi_cpa'] = fig_cpa

        with roi_commentary_col:
            roi_commentary = generate_section_commentary('roi', {
                'annual_spend': annual_spend,
                'rate_card_price': rate_card_price,
                'num_jobs': num_jobs,
                'total_clicks': total_clicks,
                'total_applies': total_applies_val,
                'cost_per_job': cost_per_job,
                'cost_per_view': cost_per_view,
                'cost_per_apply': cost_per_apply,
                'saving_pct': saving_pct,
                'roi_by_type': roi_by_type if len(roi_by_type) > 0 else None,
            })
            st.markdown(roi_commentary)
    else:
        st.metric("Jobs Advertised", f"{num_jobs:,}")
        st.info("Enter your **Annual Spend** and **Rate Card Price** in the Cost & Report Settings above to see ROI analysis.")

    st.markdown("---")

    # ===================================================================
    # SECTION 4.5: SALARY BENCHMARK BY TOP-10 OCCUPATIONS
    # ===================================================================
    # For each of the client's most-posted-with-salary roles, show the
    # market salary distribution as a histogram and overlay three reference
    # means: client, national, regional (client's HQ region). Mirrors the
    # salary-tab histogram pattern (views/salary.py:185-239) per occupation.
    st.subheader("Salary Benchmark — Your Top 10 Occupations")

    # Persistent across the conditional branches so the commentary generator
    # downstream can pick them up (None when section is skipped).
    salary_per_occ = None
    salary_client_region = None

    client_with_salary = client_df[client_df.get('has_salary_data', False) == True]

    if len(client_with_salary) == 0:
        st.info(
            "No salary data is available for this client's vacancies in the "
            "selected period — salary benchmark omitted."
        )
    else:
        occ_counts = client_with_salary['occupation'].dropna().value_counts()
        qualifying = occ_counts[occ_counts >= 5]
        top_occupations = qualifying.head(10).index.tolist()

        if len(top_occupations) == 0:
            st.info(
                "No occupations have at least 5 vacancies with salary data for "
                "this client in the selected period — salary benchmark needs "
                "≥5 priced roles per occupation to be meaningful."
            )
        else:
            # Look up client HQ region (None for multi-site / central-gov clients)
            hq_map = load_client_hq_regions()
            client_region = hq_map.get(selected_client.lower().strip())

            # Pre-build the regional market subset once. Compare on
            # lower-stripped strings so canonical/raw spelling differences
            # between primary_uk_region and client_hq_addresses don't drop
            # the line silently.
            df_regional_market = None
            if client_region and 'primary_uk_region' in df.columns:
                norm = client_region.strip().lower()
                df_regional_market = df[
                    (df.get('has_salary_data', False) == True)
                    & (df['primary_uk_region'].fillna('').str.strip().str.lower() == norm)
                ]
                if len(df_regional_market) == 0:
                    df_regional_market = None  # No samples → drop regional line

            client_color = JGP_COLORS['negative']      # red — your mean
            national_color = JGP_COLORS['amber']       # amber — national mean
            regional_color = JGP_COLORS['deep_green']  # green — regional mean

            per_occ = []
            for occ in top_occupations:
                client_occ = client_with_salary[client_with_salary['occupation'] == occ]
                client_mean = client_occ['annual_mid_salary'].mean()

                market_occ = df[(df['occupation'] == occ) & (df.get('has_salary_data', False) == True)]
                market_salaries = market_occ['annual_mid_salary'].dropna()
                national_mean = market_salaries.mean() if len(market_salaries) else np.nan

                if df_regional_market is not None:
                    reg_vals = df_regional_market[df_regional_market['occupation'] == occ]['annual_mid_salary'].dropna()
                    regional_mean = reg_vals.mean() if len(reg_vals) >= 3 else np.nan
                    regional_n = len(reg_vals)
                else:
                    regional_mean = np.nan
                    regional_n = 0

                per_occ.append({
                    'occupation': occ,
                    'client_n': len(client_occ),
                    'market_n': len(market_salaries),
                    'regional_n': regional_n,
                    'market_salaries': market_salaries,
                    'client_mean': client_mean,
                    'national_mean': national_mean,
                    'regional_mean': regional_mean,
                })

            any_regional = any(not pd.isna(p['regional_mean']) for p in per_occ)

            n_occ = len(per_occ)
            n_cols = 2
            n_rows = (n_occ + n_cols - 1) // n_cols  # ceil

            subplot_titles = [
                f"{p['occupation']} — your n={p['client_n']}, market n={p['market_n']:,}"
                for p in per_occ
            ]

            fig_salary_occ = make_subplots(
                rows=n_rows, cols=n_cols,
                subplot_titles=subplot_titles,
                vertical_spacing=0.12,
                horizontal_spacing=0.10,
            )

            for i, p in enumerate(per_occ):
                row = i // n_cols + 1
                col = i % n_cols + 1

                fig_salary_occ.add_trace(
                    go.Histogram(
                        x=p['market_salaries'],
                        nbinsx=25,
                        marker_color=JGP_COLORS['primary'],
                        opacity=0.85,
                        showlegend=False,
                        hovertemplate='Salary: £%{x:,.0f}<br>Vacancies: %{y}<extra></extra>',
                    ),
                    row=row, col=col,
                )

                if not pd.isna(p['client_mean']):
                    fig_salary_occ.add_vline(
                        x=p['client_mean'], line_width=2.5, line_color=client_color,
                        row=row, col=col,
                    )
                if not pd.isna(p['national_mean']):
                    fig_salary_occ.add_vline(
                        x=p['national_mean'], line_width=2, line_color=national_color,
                        row=row, col=col,
                    )
                if not pd.isna(p['regional_mean']):
                    fig_salary_occ.add_vline(
                        x=p['regional_mean'], line_width=2, line_color=regional_color,
                        row=row, col=col,
                    )

            # Legend traces — invisible scatters drawn once on subplot (1,1)
            # so the figure-level legend has labelled rows for each line.
            fig_salary_occ.add_trace(
                go.Scatter(
                    x=[None], y=[None], mode='lines',
                    line=dict(color=client_color, width=2.5),
                    name='Your mean',
                ),
                row=1, col=1,
            )
            fig_salary_occ.add_trace(
                go.Scatter(
                    x=[None], y=[None], mode='lines',
                    line=dict(color=national_color, width=2),
                    name='National mean',
                ),
                row=1, col=1,
            )
            if any_regional:
                fig_salary_occ.add_trace(
                    go.Scatter(
                        x=[None], y=[None], mode='lines',
                        line=dict(color=regional_color, width=2),
                        name=f"Regional mean ({client_region})",
                    ),
                    row=1, col=1,
                )

            # Two-call layout pattern avoids Python kwarg conflicts when
            # overriding template keys (see lessons.md "Spreading
            # JGP_PLOTLY_TEMPLATE['layout']").
            fig_salary_occ.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_salary_occ.update_layout(
                height=max(360, 240 * n_rows),
                showlegend=True,
                legend=dict(
                    orientation='h',
                    yanchor='bottom', y=1.04,
                    xanchor='left', x=0,
                    font=dict(size=12),
                ),
                bargap=0.05,
            )
            fig_salary_occ.update_xaxes(tickformat=',', tickprefix='£')
            fig_salary_occ.update_annotations(font_size=12)

            st.plotly_chart(fig_salary_occ, use_container_width=True)
            st.caption(CHART_EXPLAINERS['salary_by_occupation'])

            if not client_region:
                st.caption(
                    "_HQ region unavailable for this client — regional benchmark "
                    "line omitted. (Common for central-government and multi-site "
                    "bodies.)_"
                )
            elif not any_regional:
                st.caption(
                    f"_No comparable salary data found in {client_region} for "
                    f"these occupations — regional benchmark line omitted._"
                )

            report_figures['salary_by_occupation'] = fig_salary_occ

            # Expose for downstream commentary generator
            salary_per_occ = per_occ
            salary_client_region = client_region

    st.markdown("---")

    # ===================================================================
    # SECTION 5: MEDIA PERFORMANCE
    # ===================================================================
    st.subheader(f"Media Performance — {selected_client}")

    cat_stats = None  # Initialise before conditional block so it's in scope for PDF commentary
    if client_media is not None and len(client_media) > 0:
        media_vac_count = client_media['entity_id_str'].nunique() if 'entity_id_str' in client_media.columns else len(client_media)
        st.caption(f"Showing media data for **{media_vac_count:,}** vacancies belonging to {selected_client}")
        # Category-level summary
        cat_stats = client_media.groupby('source_category').agg(
            total_clicks=('clicks', 'sum'),
            total_applies=('applies', 'sum'),
            vacancy_count=('entity_id_str', 'nunique')
        ).reset_index()
        cat_stats['avg_views'] = cat_stats['total_clicks'] / cat_stats['vacancy_count']
        cat_stats['avg_applies'] = cat_stats['total_applies'] / cat_stats['vacancy_count']
        cat_stats['conversion_rate'] = (cat_stats['total_applies'] / cat_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
        cat_stats = cat_stats.sort_values('total_clicks', ascending=False)

        media_chart_col, media_commentary_col = st.columns([3, 2])

        with media_chart_col:
            fig_media = go.Figure()
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_views'],
                name='Avg. Views', orientation='h', marker_color=JGP_COLORS['primary']
            ))
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_applies'],
                name='Avg. Applies', orientation='h', marker_color=JGP_COLORS['accent'],
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_media.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_media.update_layout(
                barmode='group', height=max(350, len(cat_stats) * 40),
                title="Media Performance by Channel",
                xaxis_title="Average per Vacancy", yaxis_title="",
                legend=dict(orientation='h', y=-0.15),
                plot_bgcolor='rgba(0,0,0,0)',
            )
            st.plotly_chart(fig_media, use_container_width=True)
            st.caption(CHART_EXPLAINERS['media_performance'])
            report_figures['media'] = fig_media

            # Summary table below chart
            display_media = cat_stats[['source_category', 'vacancy_count', 'avg_views', 'avg_applies', 'conversion_rate']].copy()
            display_media.columns = ['Channel', 'Vacancies', 'Avg. Views', 'Avg. Applies', 'Conversion %']
            display_media['Avg. Views'] = display_media['Avg. Views'].round(1)
            display_media['Avg. Applies'] = display_media['Avg. Applies'].round(1)
            display_media['Conversion %'] = display_media['Conversion %'].round(1)
            st.dataframe(display_media, use_container_width=True, hide_index=True)

        with media_commentary_col:
            media_commentary = generate_section_commentary('media', {
                'cat_stats': cat_stats,
                'client_name': selected_client,
            })
            st.markdown(media_commentary)

        # Source-level detail (expandable)
        with st.expander("View by individual source"):
            media_stats = client_media.groupby(['source_category', 'source']).agg(
                total_clicks=('clicks', 'sum'),
                total_applies=('applies', 'sum'),
                vacancy_count=('entity_id_str', 'nunique')
            ).reset_index()
            media_stats['avg_views'] = media_stats['total_clicks'] / media_stats['vacancy_count']
            media_stats['avg_applies'] = media_stats['total_applies'] / media_stats['vacancy_count']
            media_stats['conversion_rate'] = (media_stats['total_applies'] / media_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
            media_stats = media_stats.sort_values('total_clicks', ascending=False)
            detail_media = media_stats[['source_category', 'source', 'vacancy_count', 'avg_views', 'avg_applies', 'conversion_rate']].copy()
            detail_media.columns = ['Channel', 'Source', 'Vacancies', 'Avg. Views', 'Avg. Applies', 'Conversion %']
            detail_media['Avg. Views'] = detail_media['Avg. Views'].round(1)
            detail_media['Avg. Applies'] = detail_media['Avg. Applies'].round(1)
            detail_media['Conversion %'] = detail_media['Conversion %'].round(1)
            st.dataframe(detail_media, use_container_width=True, hide_index=True)
    else:
        st.info("Media source data not available. Run the `dashboard_media_summary` BigQuery table creation to enable this section.")

    st.markdown("---")

    # ===================================================================
    # POWERPOINT EXPORT
    # ===================================================================
    st.subheader("Export Report")

    # --- Compute additional stats needed for PPTX template ---

    # Slide 2: Top-right quadrant % (vacancies above benchmark on BOTH views and applies)
    benchmarkable_df = scatter_df[scatter_df['category'] == 'Benchmarkable']
    if len(benchmarkable_df) > 0:
        top_quadrant_count = len(benchmarkable_df[
            (benchmarkable_df['views_diff_pct'] > 0) & (benchmarkable_df['applies_diff_pct'] > 0)
        ])
        top_quadrant_pct = (top_quadrant_count / len(benchmarkable_df)) * 100
    else:
        top_quadrant_pct = 0

    # Slide 2: Strongest job category (highest combined diff score, benchmarkable only)
    if len(benchmarkable_df) > 0:
        category_scores = benchmarkable_df.groupby('occupation').agg(
            combined_score=('views_diff_pct', lambda s: s.mean() + benchmarkable_df.loc[s.index, 'applies_diff_pct'].mean()),
            count=('views_diff_pct', 'count')
        ).reset_index()
        # Need at least 2 vacancies in occupation for the category to be considered "strongest"
        category_scores = category_scores[category_scores['count'] >= 2]
        if len(category_scores) > 0:
            top_category = category_scores.sort_values('combined_score', ascending=False).iloc[0]['occupation']
        else:
            top_category = benchmarkable_df.iloc[0]['occupation'] if len(benchmarkable_df) > 0 else 'N/A'
    else:
        top_category = 'N/A'

    # Slide 5: Build the new charts (only if spend entered)
    rate_card_total_val = rate_card_price * num_jobs if annual_spend > 0 else 0
    saving_pct_val = ((rate_card_total_val - annual_spend) / rate_card_total_val * 100) if rate_card_total_val > 0 else 0

    if annual_spend > 0:
        # Stacked bar: Your Spend (bottom) + Saving (top) = Rate Card Total
        saving_amount = max(rate_card_total_val - annual_spend, 0)
        fig_spend_stack = go.Figure()
        fig_spend_stack.add_trace(go.Bar(
            x=['Total Value'],
            y=[annual_spend],
            name='Your Spend',
            marker_color=JGP_COLORS['primary'],
            text=[f"£{annual_spend:,.0f}"],
            textposition='inside',
            textfont=dict(color='white', size=14),
            width=0.7,
        ))
        fig_spend_stack.add_trace(go.Bar(
            x=['Total Value'],
            y=[saving_amount],
            name='Saving vs Rate Card',
            marker_color=JGP_COLORS['accent'],
            text=[f"£{saving_amount:,.0f}"],
            textposition='inside',
            textfont=dict(color=JGP_COLORS['deep_blue'], size=14),
            width=0.7,
        ))
        fig_spend_stack.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
        fig_spend_stack.update_layout(
            barmode='stack',
            title=f"Your Spend vs Rate Card Value (Saving: {saving_pct_val:.0f}%)",
            yaxis_title="GBP",
            height=400,
            bargap=0.05,
            showlegend=True,
            plot_bgcolor='rgba(0,0,0,0)',
            legend=dict(orientation='h', y=-0.15),
        )
        report_figures['spend_vs_ratecard'] = fig_spend_stack

    # Cost per apply by occupation chart (always built when spend > 0)
    cpa_by_occ_fig = None
    roi_by_type_full = None
    if annual_spend > 0:
        roi_by_type_full = client_df.groupby('occupation').agg(
            total_applies=('applies', 'sum'),
            job_count=('clicks', 'count')
        ).reset_index()
        roi_by_type_full = roi_by_type_full[roi_by_type_full['total_applies'] > 0]
        if len(roi_by_type_full) > 0:
            roi_by_type_full['cost_allocated'] = annual_spend * (roi_by_type_full['job_count'] / roi_by_type_full['job_count'].sum())
            roi_by_type_full['cost_per_apply'] = roi_by_type_full['cost_allocated'] / roi_by_type_full['total_applies']
            roi_by_type_full = roi_by_type_full.sort_values('cost_per_apply', ascending=True)

            cpa_by_occ_fig = go.Figure()
            cpa_by_occ_fig.add_trace(go.Bar(
                y=roi_by_type_full['occupation'],
                x=roi_by_type_full['cost_per_apply'],
                orientation='h',
                marker_color=JGP_COLORS['amber'],
                text=roi_by_type_full['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            cpa_by_occ_fig.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            cpa_by_occ_fig.update_layout(
                title="Cost per Apply by Occupation",
                height=max(350, len(roi_by_type_full) * 32),
                xaxis_title="Cost per Apply (GBP)",
                yaxis_title="",
                plot_bgcolor='rgba(0,0,0,0)',
            )
            report_figures['cost_per_app_by_occupation'] = cpa_by_occ_fig

    # --- Build structured commentary for PPTX template ---
    scatter_struct = generate_section_commentary_structured('benchmark_scatter', {
        'total_count': len(scatter_df),
        'benchmarkable_count': len(scatter_df[scatter_df['category'] == 'Benchmarkable']),
        'zero_applies_count': len(scatter_df[scatter_df['category'].str.startswith('Zero')]),
        'no_benchmark_count': len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)']),
        'top_performers': top_performers,
        'client_name': selected_client,
    })
    average_struct = generate_section_commentary_structured('benchmark_average', {
        'client_avg_clicks': client_avg_clicks,
        'benchmark_avg_clicks': benchmark_avg_clicks,
        'client_avg_applies': client_avg_applies,
        'benchmark_avg_applies': benchmark_avg_applies,
        'num_jobs': len(client_df),
        'client_name': selected_client,
    })
    postings_struct = generate_section_commentary_structured('postings', {
        'total_jobs': len(client_df),
        'total_applies': int(client_df['applies'].sum()),
        'by_type': by_type,
        'client_name': selected_client,
    })
    roi_struct = generate_section_commentary_structured('roi', {
        'annual_spend': annual_spend,
        'rate_card_total': rate_card_total_val,
        'num_jobs': num_jobs,
        'cost_per_apply': annual_spend / total_applies_val if total_applies_val > 0 else 0,
        'saving_pct': saving_pct_val,
        'roi_by_type': roi_by_type_full if (annual_spend > 0 and roi_by_type_full is not None and len(roi_by_type_full) > 0) else None,
        'client_name': selected_client,
    })
    media_struct = generate_section_commentary_structured('media', {
        'cat_stats': cat_stats,
        'client_name': selected_client,
    })
    salary_struct = generate_section_commentary_structured('salary', {
        'per_occ': salary_per_occ,
        'client_name': selected_client,
        'client_region': salary_client_region,
    })

    # --- Build report_metrics dict (matches template tag names) ---
    report_metrics = {
        # Slide 1
        'client_name': selected_client,
        'PERIOD_START': str(report_start),
        'PERIOD_END': str(report_end),

        # Slide 2 stats
        'stat_total_jobs': f"{num_jobs:,}",
        'stat_top_quadrant_pct': f"{top_quadrant_pct:.0f}",
        'stat_top_category': top_category,

        # Slide 2 commentary
        'commentary_benchmark_intro': scatter_struct['intro'],
        'commentary_benchmark_point_1': scatter_struct['point_1'],
        'commentary_benchmark_point_2': scatter_struct['point_2'],
        'commentary_benchmark_point_3': scatter_struct['point_3'],

        # Slide 3 stats
        'stat_benchmark_average_views': f"{benchmark_avg_clicks:,.0f}",
        'stat_your_jobs_average_views': f"{client_avg_clicks:,.0f}",
        'stat_benchmark_average_applies': f"{benchmark_avg_applies:,.1f}",
        'stat_your_jobs_average_applies': f"{client_avg_applies:,.1f}",

        # Slide 3 commentary
        'commentary_average_intro': average_struct['intro'],
        'commentary_average_point_1': average_struct['point_1'],
        'commentary_average_point_2': average_struct['point_2'],

        # Slide 4 commentary
        'commentary_postings_intro': postings_struct['intro'],
        'commentary_postings_point_1': postings_struct['point_1'],
        'commentary_postings_point_2': postings_struct['point_2'],

        # Slide 5 stats (ROI)
        'stat_cost_per_job': f"£{annual_spend / num_jobs:,.2f}" if (annual_spend > 0 and num_jobs > 0) else "—",
        'stat_cost_per_view': f"£{annual_spend / total_clicks:,.2f}" if (annual_spend > 0 and total_clicks > 0) else "—",
        'stat_cost_per_apply': f"£{annual_spend / total_applies_val:,.2f}" if (annual_spend > 0 and total_applies_val > 0) else "—",

        # Slide 5 commentary
        'commentary_roi_intro': roi_struct['intro'],
        'commentary_roi_point_1': roi_struct['point_1'],
        'commentary_roi_point_2': roi_struct['point_2'],

        # Slide 6 commentary (Salary Benchmark)
        'commentary_salary_intro': salary_struct['intro'],
        'commentary_salary_point_1': salary_struct['point_1'],
        'commentary_salary_point_2': salary_struct['point_2'],
        'commentary_salary_point_3': salary_struct['point_3'],

        # Slide 7 commentary
        'commentary_media_intro': media_struct['intro'],
        'commentary_media_point_1': media_struct['point_1'],
        'commentary_media_point_2': media_struct['point_2'],
        'commentary_media_point_3': media_struct['point_3'],

        # Static chart explainers — sourced from module-level CHART_EXPLAINERS
        # so on-screen captions and PPTX placeholders share one source of truth.
        **{f'chart_explainer_{k}': v for k, v in CHART_EXPLAINERS.items()},

        # Slide 7 contact
        'contact_name': contact_name or 'Your Account Manager',
        'contact_title': contact_title or 'Account Director',
        'contact_email': contact_email or 'team@jobsgopublic.com',
        'contact_phone': contact_phone or '020 7427 8250',
    }

    # --- Map template chart tags to figures ---
    pptx_figures = {
        'benchmark_scatter': report_figures.get('scatter'),
        'benchmark_average': report_figures.get('benchmark_combined'),
        'postings_by_type': report_figures.get('postings'),
        'spend_vs_ratecard': report_figures.get('spend_vs_ratecard'),
        'cost_per_app_by_occupation': report_figures.get('cost_per_app_by_occupation'),
        'salary_by_occupation': report_figures.get('salary_by_occupation'),
        'media_performance': report_figures.get('media'),
    }

    # --- Generate PPTX and offer download ---
    template_path = 'Renewals.pptx'
    try:
        pptx_bytes = generate_client_report_pptx(report_metrics, pptx_figures, template_path)
        st.download_button(
            "Download PowerPoint Report",
            data=pptx_bytes,
            file_name=f"advertising_report_{selected_client.replace(' ', '_')}_{report_start}_{report_end}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary"
        )
        st.caption("Download the PowerPoint, edit any commentary you'd like, then File → Export → PDF for the final client copy.")
    except FileNotFoundError:
        st.error(f"Template file not found at `{template_path}`. Make sure `Renewals.pptx` is in the project root.")
    except Exception as e:
        st.warning("PowerPoint generation requires `python-pptx` and `kaleido`.")
        st.caption(f"Error: {type(e).__name__}: {e}")


def generate_client_report_pptx(metrics, figures, template_path):
    """Generate a PowerPoint report by filling a branded template.

    Args:
        metrics: dict mapping placeholder names (without {{}}) to string values
        figures: dict mapping chart slot names (e.g. 'benchmark_scatter') to plotly figures
        template_path: path to the .pptx template file

    Returns: bytes of the populated .pptx file.

    Replaces text placeholders like {{tag_name}} with metrics[tag_name].
    Replaces chart placeholders like {{chart:slot_name}} with PNG images of figures[slot_name].
    """
    prs = Presentation(template_path)

    # --- Helper: render a Plotly figure to PNG bytes at a target aspect ratio ---
    def _fig_to_png(fig, slot_width_emu=None, slot_height_emu=None):
        """Render a Plotly figure to white-background PNG bytes.

        If slot dimensions are passed (the placeholder's .width/.height in
        EMU), the PNG is rendered at that exact aspect ratio so PowerPoint
        doesn't squash the bitmap when fitting it. Otherwise default to 16:9.

        Pixel target is 1500 on the longest side, then kaleido renders at
        scale=4 — effectively very high DPI, sharp at any zoom and when
        exported to PDF. Y/X axes use automargin so tick labels never clip.
        """
        if fig is None:
            return None
        try:
            # Match the placeholder aspect ratio but keep canvas dimensions
            # constant across charts. Plotly font sizes are pixels relative to
            # canvas, so a uniform canvas keeps text/legend/axes the same
            # relative size on every slide. Crispness comes from `scale` (a
            # pure DPI multiplier), not from a bigger canvas.
            TARGET_LONG_PX = 1800
            if slot_width_emu and slot_height_emu:
                ratio = slot_width_emu / slot_height_emu
            else:
                ratio = 16 / 9
            if ratio >= 1:
                width = TARGET_LONG_PX
                height = int(round(TARGET_LONG_PX / ratio))
            else:
                height = TARGET_LONG_PX
                width = int(round(TARGET_LONG_PX * ratio))

            fig_export = go.Figure(fig.to_dict())
            fig_export.update_layout(
                paper_bgcolor='white',
                plot_bgcolor='white',
                font=dict(size=18),
                # pad=15 puts visible space between tick labels and the plot.
                margin=dict(l=120, r=40, t=60, b=80, pad=15),
                # Override the JGP template's small (12pt) legend font for export.
                legend=dict(font=dict(size=18)),
            )
            fig_export.update_xaxes(automargin=True, title_font=dict(size=22))
            fig_export.update_yaxes(automargin=True, title_font=dict(size=22))
            # Force-set bar value-label size and disable auto-shrink. Plotly
            # defaults constraintext='both' which silently shrinks text to fit
            # inside the bar — that's why update_traces alone wasn't visibly
            # bumping the dense postings chart and the narrow spend column.
            for trace in fig_export.data:
                if trace.type != 'bar':
                    continue
                preserved_color = None
                if trace.textfont is not None and trace.textfont.color is not None:
                    preserved_color = trace.textfont.color
                trace.textfont = ({'size': 22, 'color': preserved_color}
                                  if preserved_color else {'size': 22})
                trace.constraintext = 'none'

            # scale=6 → 10800 px on the long side. Maximum crispness; render
            # time is a few seconds per chart, accepted for renewals reports.
            return fig_export.to_image(format='png', width=width, height=height, scale=6)
        except Exception:
            return None

    # --- Helper: replace text in a single shape's text frame, preserving formatting ---
    def _replace_text_in_shape(shape, replacements):
        """Walk runs/paragraphs and replace {{tag}} occurrences. Handles tags split across runs."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            # First try simple per-run replacement (works when tag is fully in one run)
            for run in paragraph.runs:
                text = run.text
                if '{{' in text:
                    for tag, val in replacements.items():
                        placeholder = '{{' + tag + '}}'
                        if placeholder in text:
                            text = text.replace(placeholder, str(val))
                    run.text = text

            # Fallback: if a tag spans multiple runs, the above won't catch it.
            # Concatenate paragraph text, replace, and put it all in the first run.
            full_text = ''.join(r.text for r in paragraph.runs)
            if '{{' in full_text and any('{{' + tag + '}}' in full_text for tag in replacements):
                new_text = full_text
                for tag, val in replacements.items():
                    new_text = new_text.replace('{{' + tag + '}}', str(val))
                if new_text != full_text:
                    if paragraph.runs:
                        paragraph.runs[0].text = new_text
                        for r in paragraph.runs[1:]:
                            r.text = ''

    # --- Step 1: Find all chart placeholders, capture position/size, queue for replacement ---
    # We do this before text replacement so we can find {{chart:xxx}} markers.
    chart_replacements = []  # list of (slide, shape, slot_name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                m = re.search(r'\{\{chart:([^}]+)\}\}', text)
                if m:
                    chart_replacements.append((slide, shape, m.group(1)))

    # --- Step 2: Replace text placeholders on every shape across all slides ---
    # `slide_number` and `total_slides` are computed per-slide so a single
    # template footer like `{{slide_number}} / {{total_slides}}` renders as
    # "1 / 9" on slide 1, "2 / 9" on slide 2, etc.
    text_replacements = {k: v for k, v in metrics.items()}
    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        per_slide = {
            **text_replacements,
            'slide_number': idx,
            'total_slides': total_slides,
        }
        for shape in slide.shapes:
            _replace_text_in_shape(shape, per_slide)

    # --- Step 3: Replace chart placeholders with images ---
    for slide, shape, slot_name in chart_replacements:
        fig = figures.get(slot_name)
        # Capture original geometry before deleting the shape
        left, top, width, height = shape.left, shape.top, shape.width, shape.height

        # Remove the placeholder shape
        sp = shape._element
        sp.getparent().remove(sp)

        if fig is not None:
            png_bytes = _fig_to_png(fig, slot_width_emu=width, slot_height_emu=height)
            if png_bytes:
                slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width, height=height)
            else:
                # Failed to render — leave a small note in place
                txt_box = slide.shapes.add_textbox(left, top, width, height)
                txt_box.text_frame.text = '[Chart unavailable]'
        # If fig is None, just remove the placeholder silently

    # --- Step 4: Clean up any remaining unreplaced {{tags}} (set to empty so they don't show) ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join(r.text for r in paragraph.runs)
                    if re.search(r'\{\{[^}]+\}\}', full_text):
                        cleaned = re.sub(r'\{\{[^}]+\}\}', '', full_text)
                        if paragraph.runs:
                            paragraph.runs[0].text = cleaned
                            for r in paragraph.runs[1:]:
                                r.text = ''

    # --- Step 5: Output to bytes ---
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
