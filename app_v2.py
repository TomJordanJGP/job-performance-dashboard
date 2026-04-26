import streamlit as st
import pandas as pd
import numpy as np
from google.oauth2.service_account import Credentials
from google.cloud import bigquery
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
from data.processing import apply_media_categories, process_salary_columns
from views.salary import render_salary

# Page configuration
st.set_page_config(
    page_title="Job Performance Dashboard (v2 Dev)",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# BigQuery configuration
BQ_PROJECT_ID = "site-monitoring-421401"
BQ_DATASET_ID = "job_data_export"
BQ_TABLE_ID = "dashboard_vacancy_summary"
BQ_DAILY_TOTALS_TABLE_ID = "dashboard_daily_totals"
BQ_MEDIA_SUMMARY_TABLE_ID = "dashboard_media_summary"
BQ_REGION_SUMMARY_TABLE_ID = "dashboard_vacancy_region_summary"

SCOPES = [
    'https://www.googleapis.com/auth/spreadsheets.readonly',
    'https://www.googleapis.com/auth/bigquery',
]

# ============================================================================
# DATA LOADING FUNCTIONS
# ============================================================================

@st.cache_resource(ttl=None)
def get_bigquery_client():
    """Initialize and cache the BigQuery client."""
    import os

    # Get the directory where this script is located
    script_dir = os.path.dirname(os.path.abspath(__file__))
    service_account_path = os.path.join(script_dir, 'service_account.json')

    try:
        # Try Streamlit secrets first (for cloud deployment)
        use_secrets = False
        try:
            if hasattr(st, 'secrets') and 'gcp_service_account' in st.secrets:
                use_secrets = True
        except Exception as e:
            use_secrets = False

        if use_secrets:
            creds = Credentials.from_service_account_info(
                st.secrets['gcp_service_account'],
                scopes=SCOPES
            )
        else:
            # Fall back to local file (for local development)
            pass

            # Check if file exists before trying to use it
            if not os.path.exists(service_account_path):
                st.error(f"❌ No authentication found!")
                st.error(f"Local file does not exist at: {service_account_path}")
                st.error("Please either:")
                st.error("1. Add secrets to Streamlit Cloud (Settings → Secrets), OR")
                st.error("2. Add service_account.json file to the app directory")
                st.stop()

            creds = Credentials.from_service_account_file(
                service_account_path,
                scopes=SCOPES
            )

        client = bigquery.Client(credentials=creds, project=BQ_PROJECT_ID)
        return client
    except FileNotFoundError as e:
        st.error(f"⚠️ Service account credentials not found at: {service_account_path}")
        st.error("Please add them to Streamlit secrets or place service_account.json in the app directory")
        st.code(f"Expected location: {service_account_path}")
        st.code(f"Error: {repr(e)}")
        st.stop()
    except Exception as e:
        st.error(f"❌ Unexpected error initializing BigQuery client: {type(e).__name__}")
        st.error(f"Error message: {str(e)}")
        st.code(f"Attempted to load from: {service_account_path}")
        st.code(f"Full error: {repr(e)}")
        st.stop()

def _ensure_media_summary_table(client):
    """Create dashboard_media_summary if it doesn't exist, then load it.

    This runs once — after that the daily scheduled query keeps it fresh
    (once create_aggregated_tables.sql is updated in BigQuery Console).
    """
    create_sql = f"""
    CREATE TABLE IF NOT EXISTS `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_MEDIA_SUMMARY_TABLE_ID}`
    AS
    SELECT
      entity_id_str,
      importer_ID,
      ANY_VALUE(importer_name) as importer_name,
      source,
      medium,
      campaign,
      COUNTIF(event_name = 'job_visit') as clicks,
      COUNTIF(event_name = 'job_apply_start') as applies
    FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.job_performance_enriched`
    WHERE event_name IN ('job_visit', 'job_apply_start')
    GROUP BY entity_id_str, importer_ID, source, medium, campaign
    """
    try:
        job = client.query(create_sql)
        job.result()

        read_sql = f"""
        SELECT entity_id_str, importer_ID, importer_name,
               source, medium, campaign, clicks, applies
        FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_MEDIA_SUMMARY_TABLE_ID}`
        """
        df = client.query(read_sql).to_dataframe(create_bqstorage_client=False)
        return df
    except Exception:
        return None


@st.cache_data(ttl=14400)  # Cache for 4 hours (data refreshes daily)
def load_all_data(days_back=30, sample_size=None):
    """Load vacancy summary, daily totals, and media summary.

    Returns three DataFrames. Media query is optional — returns None if unavailable.

    Args:
        days_back: Number of days to look back
        sample_size: If set, limit vacancy rows (for testing)
    """
    try:
        client = get_bigquery_client()

        cutoff_date = (datetime.now() - timedelta(days=days_back)).strftime('%Y-%m-%d')

        limit_clause = f"LIMIT {sample_size}" if sample_size else ""

        # Core fields always present in dashboard_vacancy_summary
        core_fields = """
            entity_id_str,
            first_event_date,
            last_event_date,
            clicks,
            applies,
            title,
            organization_name,
            uk_regions,
            primary_uk_region,
            occupational_fields,
            importer_ID,
            importer_name,
            workflow_state,
            upgrades,
            start_date,
            end_date,
            category,
            contract_type,
            employment_type"""

        # Salary fields (added after running updated create_aggregated_tables.sql)
        salary_fields = """,
            min_salary,
            max_salary,
            currency_code,
            salary_free_text,
            salary_exact,
            salary_unit"""

        vacancy_query = f"""
        SELECT {core_fields}{salary_fields}
        FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_TABLE_ID}`
        WHERE last_event_date >= '{cutoff_date}'
        {limit_clause}
        """

        daily_query = f"""
        SELECT
            event_date,
            clicks,
            applies,
            active_vacancies
        FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_DAILY_TOTALS_TABLE_ID}`
        WHERE event_date >= '{cutoff_date}'
        ORDER BY event_date
        """

        media_query = f"""
        SELECT
            entity_id_str,
            importer_ID,
            importer_name,
            source,
            medium,
            campaign,
            clicks,
            applies
        FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_MEDIA_SUMMARY_TABLE_ID}`
        """

        # Try with salary fields; fall back to core-only if table not yet updated
        try:
            vacancy_job = client.query(vacancy_query)
            vacancy_job.result()
        except Exception:
            vacancy_query = f"""
            SELECT {core_fields}
            FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_TABLE_ID}`
            WHERE last_event_date >= '{cutoff_date}'
            {limit_clause}
            """
            vacancy_job = client.query(vacancy_query)
            vacancy_job.result()

        daily_job = client.query(daily_query)
        media_job = None
        try:
            media_job = client.query(media_query)
        except Exception:
            pass

        # Wait for daily to complete (vacancy already resolved above)
        daily_job.result()

        vacancy_df = vacancy_job.to_dataframe(create_bqstorage_client=False)
        daily_df = daily_job.to_dataframe(create_bqstorage_client=False)

        # Try to load media data; auto-create the table if it doesn't exist
        media_df = None
        if media_job is not None:
            try:
                media_job.result()
                media_df = media_job.to_dataframe(create_bqstorage_client=False)
            except Exception:
                media_df = _ensure_media_summary_table(client)

        # Region-exploded summary (one row per vacancy per region)
        region_df = None
        try:
            region_query = f"""
            SELECT *
            FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.{BQ_REGION_SUMMARY_TABLE_ID}`
            WHERE last_event_date >= '{cutoff_date}'
            """
            region_job = client.query(region_query)
            region_job.result()
            region_df = region_job.to_dataframe(create_bqstorage_client=False)
        except Exception:
            pass  # Table may not exist yet

        return vacancy_df, daily_df, media_df, region_df
    except Exception as e:
        st.error(f"❌ Error loading data: {str(e)}")
        st.markdown("""
        **Troubleshooting:**
        - Check BigQuery tables exist: `dashboard_vacancy_summary` and `dashboard_daily_totals`
        - Verify service account has `bigquery.jobs.create` permission
        - Check the tables have data for the requested date range
        """)
        st.stop()

@st.cache_data(ttl=300)
def load_importer_mapping():
    """Load importer mapping from CSV file."""
    try:
        mapping_df = pd.read_csv('importer_mapping.csv', encoding='utf-8-sig')
        if 'importer_id' in mapping_df.columns and 'importer_name' in mapping_df.columns:
            mapping_df = mapping_df[mapping_df['importer_id'].notna()]
            mapping_df = mapping_df[mapping_df['importer_id'].astype(str).str.strip() != '']
            # Create mapping with string keys (stripped)
            importer_mapping = dict(zip(
                mapping_df['importer_id'].astype(str).str.strip(),
                mapping_df['importer_name'].str.strip()
            ))
            return importer_mapping
        return {}
    except Exception as e:
        st.error(f"Error loading importer mapping: {e}")
        return {}

@st.cache_data(ttl=14400)
def load_launch_timing_data(days_back=365):
    """Load per-vacancy per-day event counts from the enriched table for launch timing analysis."""
    try:
        client = get_bigquery_client()
        cutoff_date = (datetime.now() - timedelta(days=days_back)).strftime('%Y-%m-%d')

        query = f"""
        WITH vacancy_start AS (
            SELECT
                entity_id_str,
                MIN(event_date_parsed) as first_event_date,
                ANY_VALUE(occupational_fields) as occupational_fields
            FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.job_performance_enriched`
            WHERE event_date_parsed >= '{cutoff_date}'
            GROUP BY entity_id_str
        ),
        daily_events AS (
            SELECT
                e.entity_id_str,
                e.event_date_parsed,
                COUNTIF(e.event_name = 'job_visit') as clicks,
                COUNTIF(e.event_name = 'job_apply_start') as applies
            FROM `{BQ_PROJECT_ID}.{BQ_DATASET_ID}.job_performance_enriched` e
            WHERE e.event_date_parsed >= '{cutoff_date}'
            GROUP BY e.entity_id_str, e.event_date_parsed
        )
        SELECT
            d.entity_id_str,
            d.event_date_parsed,
            d.clicks,
            d.applies,
            DATE_DIFF(d.event_date_parsed, v.first_event_date, DAY) as day_offset,
            EXTRACT(DAYOFWEEK FROM v.first_event_date) as launch_dow,
            EXTRACT(DAYOFWEEK FROM d.event_date_parsed) as event_dow,
            v.occupational_fields
        FROM daily_events d
        JOIN vacancy_start v ON d.entity_id_str = v.entity_id_str
        WHERE DATE_DIFF(d.event_date_parsed, v.first_event_date, DAY) BETWEEN 0 AND 30
        ORDER BY d.entity_id_str, d.event_date_parsed
        """

        job = client.query(query)
        job.result()
        return job.to_dataframe(create_bqstorage_client=False)
    except Exception as e:
        st.error(f"Error loading launch timing data: {str(e)}")
        return pd.DataFrame()


# ============================================================================
# DATA PROCESSING FUNCTIONS
# ============================================================================

def apply_importer_mapping(df, mapping):
    """Apply importer ID to name mapping. Uses BigQuery importer_name as primary,
    falls back to CSV mapping for any NULLs."""
    if 'importer_name' in df.columns:
        # importer_name already comes from BigQuery; fill gaps with CSV mapping
        df = df.copy()
        if mapping and 'importer_ID' in df.columns:
            df['importer_id_str'] = df['importer_ID'].astype(str).str.strip()
            mask = df['importer_name'].isna() | (df['importer_name'].astype(str).str.strip() == '')
            df.loc[mask, 'importer_name'] = df.loc[mask, 'importer_id_str'].map(mapping)
        df['importer_name'] = df['importer_name'].fillna('Unknown')
        return df

    if 'importer_ID' not in df.columns:
        df['importer_name'] = 'Unknown'
        return df

    # Fallback: no importer_name column — use CSV mapping
    df = df.copy()
    df['importer_id_str'] = df['importer_ID'].astype(str).str.strip()

    if mapping:
        df['importer_name'] = df['importer_id_str'].map(mapping)
        df['importer_name'] = df['importer_name'].fillna('ID: ' + df['importer_id_str'])
    else:
        df['importer_name'] = 'ID: ' + df['importer_id_str']

    return df

def parse_upgrades(df):
    """Parse upgrades column and create individual upgrade columns."""
    if 'upgrades' not in df.columns:
        return df

    # Create a copy to avoid modifying original
    df = df.copy()

    # Extract all unique upgrade types
    all_upgrades = set()
    for upgrades_str in df['upgrades'].dropna():
        if pd.notna(upgrades_str) and upgrades_str.strip():
            upgrades_list = [u.strip() for u in str(upgrades_str).split('|')]
            all_upgrades.update(upgrades_list)

    # Store parsed upgrades as a list
    df['upgrades_list'] = df['upgrades'].apply(lambda x:
        [u.strip() for u in str(x).split('|')] if pd.notna(x) and str(x).strip() else []
    )

    return df

def prepare_enriched_data(df):
    """Prepare vacancy summary data by renaming columns for dashboard compatibility."""
    df = df.copy()

    # The vacancy summary table already uses clean names for most columns.
    # Only entity_id_str needs renaming.
    column_mapping = {
        'entity_id_str': 'entity_id',
    }

    # Only rename columns that exist
    existing_renames = {k: v for k, v in column_mapping.items() if k in df.columns}
    df = df.rename(columns=existing_renames)

    return df

def add_occupation_column(df):
    """Extract occupation field from occupational_fields column."""
    if 'occupational_fields' in df.columns:
        df['occupation'] = df['occupational_fields'].apply(lambda x:
            str(x).split('|')[0].strip() if pd.notna(x) and str(x).strip() else 'Unknown'
        )
    else:
        df['occupation'] = 'Unknown'

    return df

def parse_dates_in_jobiqo(df):
    """Parse date columns from vacancy summary data."""
    if 'first_event_date' in df.columns:
        df['first_event_date'] = pd.to_datetime(df['first_event_date'], errors='coerce', utc=True).dt.tz_localize(None)
    if 'last_event_date' in df.columns:
        df['last_event_date'] = pd.to_datetime(df['last_event_date'], errors='coerce', utc=True).dt.tz_localize(None)
    if 'start_date' in df.columns:
        df['start_date'] = pd.to_datetime(df['start_date'], errors='coerce', utc=True).dt.tz_localize(None)
    if 'end_date' in df.columns:
        df['end_date'] = pd.to_datetime(df['end_date'], errors='coerce', utc=True).dt.tz_localize(None)
    return df

# ============================================================================
# FILTER FUNCTIONS
# ============================================================================

def create_filter_panel(df, key_prefix, default_months=6):
    """Create a reusable filter panel for all tabs with compact 3-column layout."""
    st.subheader("🔍 Filters")

    filters = {}

    # Row 1: Date Range, Importer, Company
    col1, col2, col3 = st.columns(3)

    with col1:
        # Date Range Filter (uses last_event_date for range boundaries)
        date_col = 'last_event_date' if 'last_event_date' in df.columns else None
        if date_col and pd.api.types.is_datetime64_any_dtype(df[date_col]):
            # Use first_event_date min and last_event_date max for the full range
            min_date = df['first_event_date'].dropna().min() if 'first_event_date' in df.columns else df[date_col].dropna().min()
            max_date = df[date_col].dropna().max()
            if pd.isna(min_date) or pd.isna(max_date):
                min_date = datetime.now().date()
                max_date = datetime.now().date()
            else:
                min_date = min_date.date() if hasattr(min_date, 'date') else min_date
                max_date = max_date.date() if hasattr(max_date, 'date') else max_date
            default_start = min_date

            filters['date_range'] = st.date_input(
                "Date Range",
                [default_start, max_date],
                min_value=min_date,
                max_value=max_date,
                key=f'{key_prefix}_date'
            )

    with col2:
        # Importer Filter
        if 'importer_name' in df.columns:
            importers = sorted(df['importer_name'].dropna().unique())
            filters['importer'] = st.multiselect(
                "Importer",
                importers,
                key=f'{key_prefix}_importer'
            )

    with col3:
        # Company Filter
        if 'organization_name' in df.columns:
            companies = sorted(df['organization_name'].dropna().unique())
            filters['company'] = st.multiselect(
                "Company",
                companies,
                key=f'{key_prefix}_company'
            )

    # Row 2: Region, Occupation, Upgrades
    col1, col2, col3 = st.columns(3)

    with col1:
        # Region Filter — use pre-exploded region_df when available
        _rdf = st.session_state.get('_region_df')
        if _rdf is not None and 'uk_region' in _rdf.columns:
            regions = sorted(_rdf['uk_region'].dropna().unique())
        elif 'uk_regions' in df.columns:
            all_regions = set()
            for regions_str in df['uk_regions'].dropna():
                for r in str(regions_str).split(' | '):
                    r = r.strip()
                    if r:
                        all_regions.add(r)
            regions = sorted(all_regions)
        else:
            regions = []
        if regions:
            filters['region'] = st.multiselect(
                "Region",
                regions,
                key=f'{key_prefix}_region'
            )

    with col2:
        # Occupation Filter
        if 'occupation' in df.columns:
            occupations = sorted(df['occupation'].dropna().unique())
            filters['occupation'] = st.multiselect(
                "Occupation",
                occupations,
                key=f'{key_prefix}_occupation'
            )

    with col3:
        # Upgrades Filter
        if 'upgrades_list' in df.columns:
            all_upgrades = set()
            for upgrades in df['upgrades_list']:
                all_upgrades.update(upgrades)
            upgrade_options = sorted(list(all_upgrades))
            filters['upgrades'] = st.multiselect(
                "Upgrades",
                upgrade_options,
                key=f'{key_prefix}_upgrades'
            )

    # Row 3: Job Title Search and Apply Button
    col1, col2, col3 = st.columns([2, 1, 1])

    with col1:
        # Job Title Search
        filters['job_title'] = st.text_input(
            "Job Title (search)",
            key=f'{key_prefix}_title',
            placeholder="e.g., Housing Director"
        )

    with col2:
        # Apply button
        st.write("")  # Spacer to align button
        apply_clicked = st.button(
            "🔄 Apply Filters",
            key=f'{key_prefix}_apply',
            type="primary",
            width='stretch'
        )

    st.markdown("---")

    return filters, apply_clicked

def apply_filters_to_data(df, filters):
    """Apply filter selections to dataframe."""
    # Handle None filters (no filters applied yet)
    if filters is None:
        return df.copy()

    filtered = df.copy()

    # Date Range Filter (show vacancies whose event date range overlaps the selected range)
    if filters.get('date_range') and len(filters['date_range']) == 2:
        start_date, end_date = filters['date_range']

        # A vacancy is in range if first_event_date <= end_date AND last_event_date >= start_date
        if 'first_event_date' in filtered.columns and 'last_event_date' in filtered.columns:
            if pd.api.types.is_datetime64_any_dtype(filtered['last_event_date']):
                filtered = filtered[
                    (filtered['first_event_date'].dt.date <= end_date) &
                    (filtered['last_event_date'].dt.date >= start_date)
                ]

    # Importer Filter
    if filters.get('importer') and 'importer_name' in filtered.columns:
        filtered = filtered[filtered['importer_name'].isin(filters['importer'])]

    # Company Filter
    if filters.get('company') and 'organization_name' in filtered.columns:
        filtered = filtered[filtered['organization_name'].isin(filters['company'])]

    # Region Filter — match vacancies where ANY region in pipe-separated uk_regions matches
    if filters.get('region') and 'uk_regions' in filtered.columns:
        selected_regions = set(filters['region'])
        mask = filtered['uk_regions'].apply(
            lambda x: bool(selected_regions & set(r.strip() for r in str(x).split(' | ')))
            if pd.notna(x) else False
        )
        filtered = filtered[mask]

    # Occupation Filter
    if filters.get('occupation') and 'occupation' in filtered.columns:
        filtered = filtered[filtered['occupation'].isin(filters['occupation'])]

    # Upgrades Filter (vacancy has ANY of the selected upgrades)
    if filters.get('upgrades') and 'upgrades_list' in filtered.columns:
        filtered = filtered[filtered['upgrades_list'].apply(
            lambda x: any(upgrade in x for upgrade in filters['upgrades'])
        )]

    # Job Title Search (case-insensitive partial match)
    if filters.get('job_title') and filters['job_title'].strip():
        if 'title' in filtered.columns:
            search_term = filters['job_title'].strip().lower()
            filtered = filtered[filtered['title'].str.lower().str.contains(search_term, na=False)]

    return filtered

# ============================================================================
# CALCULATION FUNCTIONS
# ============================================================================

def remove_outliers_iqr(data):
    """Remove outliers using IQR (Interquartile Range) method."""
    if len(data) < 4:  # Need at least 4 points for IQR
        return data

    q1 = np.percentile(data, 25)
    q3 = np.percentile(data, 75)
    iqr = q3 - q1

    lower_bound = q1 - (1.5 * iqr)
    upper_bound = q3 + (1.5 * iqr)

    return [x for x in data if lower_bound <= x <= upper_bound]


def calculate_metrics(df):
    """Calculate key metrics from pre-aggregated vacancy summary data.

    Each row in df is a vacancy with pre-computed 'clicks' and 'applies' columns.
    """
    entity_col = 'entity_id' if 'entity_id' in df.columns else df.columns[0]

    metrics = {}
    metrics['num_vacancies'] = len(df)  # Each row IS a vacancy now

    if 'clicks' in df.columns:
        metrics['total_clicks'] = int(df['clicks'].sum())
        metrics['total_applies'] = int(df['applies'].sum())
    else:
        metrics['total_clicks'] = 0
        metrics['total_applies'] = 0

    metrics['apply_click_ratio'] = (metrics['total_applies'] / metrics['total_clicks'] * 100) if metrics['total_clicks'] > 0 else 0

    # Calculate per-vacancy metrics directly from columns
    if metrics['num_vacancies'] > 0 and 'clicks' in df.columns:
        metrics['mean_clicks_per_vacancy'] = metrics['total_clicks'] / metrics['num_vacancies']
        metrics['mean_applies_per_vacancy'] = metrics['total_applies'] / metrics['num_vacancies']
        metrics['median_clicks_per_vacancy'] = float(np.median(df['clicks'].values))
        metrics['median_applies_per_vacancy'] = float(np.median(df['applies'].values))
        metrics['clicks_per_vacancy'] = metrics['mean_clicks_per_vacancy']
        metrics['applies_per_vacancy'] = metrics['mean_applies_per_vacancy']
    else:
        metrics['median_clicks_per_vacancy'] = 0
        metrics['median_applies_per_vacancy'] = 0
        metrics['mean_clicks_per_vacancy'] = 0
        metrics['mean_applies_per_vacancy'] = 0
        metrics['clicks_per_vacancy'] = 0
        metrics['applies_per_vacancy'] = 0

    return metrics

def calculate_quartile_metrics(df):
    """Calculate metrics by performance quartiles (top 25%, middle 50%, bottom 25%).

    Uses pre-aggregated clicks and applies columns directly.
    """
    if 'clicks' not in df.columns:
        return None

    if len(df) < 4:
        return None  # Need at least 4 vacancies for quartiles

    vacancy_clicks = df['clicks']
    vacancy_applies = df['applies']

    # Calculate quartile thresholds based on clicks
    q1_threshold = vacancy_clicks.quantile(0.25)
    q3_threshold = vacancy_clicks.quantile(0.75)

    # Categorize vacancies using vectorized operations
    top_25_mask = vacancy_clicks >= q3_threshold
    middle_50_mask = (vacancy_clicks >= q1_threshold) & (vacancy_clicks < q3_threshold)
    bottom_25_mask = vacancy_clicks < q1_threshold

    # Calculate metrics for each quartile
    quartiles = {}

    for name, mask in [('top_25', top_25_mask), ('middle_50', middle_50_mask), ('bottom_25', bottom_25_mask)]:
        total_clicks = int(vacancy_clicks[mask].sum())
        total_applies = int(vacancy_applies[mask].sum())
        num_vacancies = int(mask.sum())

        quartiles[name] = {
            'num_vacancies': num_vacancies,
            'total_clicks': total_clicks,
            'total_applies': total_applies,
            'apply_click_ratio': (total_applies / total_clicks * 100) if total_clicks > 0 else 0,
            'clicks_per_vacancy': total_clicks / num_vacancies if num_vacancies > 0 else 0,
            'applies_per_vacancy': total_applies / num_vacancies if num_vacancies > 0 else 0
        }

    return quartiles


def get_performance_color(value, avg_value, metric_type='ratio'):
    """Get color indicator based on performance vs average."""
    if value is None or avg_value is None or avg_value == 0:
        return "⚪"

    diff_pct = ((value - avg_value) / avg_value) * 100

    if diff_pct > 10:
        return "🟢"  # Above average
    elif diff_pct < -10:
        return "🔴"  # Below average
    else:
        return "🟡"  # Near average

# ============================================================================
# TAB 1: OVERVIEW DASHBOARD
# ============================================================================

def create_overview_tab(df, daily_totals=None):
    """Create the Overview Dashboard tab.

    Args:
        df: The vacancy summary dataframe (one row per vacancy)
        daily_totals: Pre-aggregated daily totals dataframe (one row per day)
    """
    st.header("📊 Overview Dashboard")

    # Filters in sidebar/expander
    with st.expander("🔍 Filters", expanded=True):
        filters, apply_clicked = create_filter_panel(df, 'overview')

    # Apply filters
    if apply_clicked or 'overview_filters' in st.session_state:
        if apply_clicked:
            st.session_state.overview_filters = filters
        filtered_df = apply_filters_to_data(df, st.session_state.overview_filters)
    else:
        filtered_df = df.copy()

    # Calculate metrics
    metrics = calculate_metrics(filtered_df)
    quartiles = calculate_quartile_metrics(filtered_df)

    # Debug: Show what we got
    with st.expander("🔧 Debug Info", expanded=False):
        st.write("**Available columns:**")
        st.write(filtered_df.columns.tolist())
        st.write("**Metrics calculated:**")
        st.json(metrics)
        st.write("**Sample data (first 5 rows):**")
        st.dataframe(filtered_df.head())

    # KPI Cards with Quartile Breakdown
    st.subheader("Key Performance Indicators")

    if quartiles:
        # Row 1: Overall Totals
        st.markdown("### 📊 Overall Performance")
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric("Total Vacancies", f"{metrics['num_vacancies']:,}")

        with col2:
            st.metric("Total Clicks", f"{metrics['total_clicks']:,}")

        with col3:
            st.metric("Total Applies", f"{metrics['total_applies']:,}")

        with col4:
            st.metric("Apply/Click Ratio", f"{metrics['apply_click_ratio']:.2f}%")

        st.markdown("---")

        # Row 2: Top 25% (Best Performers)
        st.markdown("### 🟢 Top 25% Performers")
        col1, col2, col3, col4, col5, col6 = st.columns(6)

        with col1:
            st.metric("Vacancies", f"{quartiles['top_25']['num_vacancies']:,}")

        with col2:
            st.metric("Total Clicks", f"{quartiles['top_25']['total_clicks']:,}")

        with col3:
            st.metric("Total Applies", f"{quartiles['top_25']['total_applies']:,}")

        with col4:
            st.metric("Apply/Click %", f"{quartiles['top_25']['apply_click_ratio']:.2f}%")

        with col5:
            st.metric("Avg Clicks/Vac", f"{quartiles['top_25']['clicks_per_vacancy']:.1f}")

        with col6:
            st.metric("Avg Applies/Vac", f"{quartiles['top_25']['applies_per_vacancy']:.2f}")

        st.markdown("---")

        # Row 3: Middle 50% (Average Performers)
        st.markdown("### 🟡 Middle 50% Performers")
        col1, col2, col3, col4, col5, col6 = st.columns(6)

        with col1:
            st.metric("Vacancies", f"{quartiles['middle_50']['num_vacancies']:,}")

        with col2:
            st.metric("Total Clicks", f"{quartiles['middle_50']['total_clicks']:,}")

        with col3:
            st.metric("Total Applies", f"{quartiles['middle_50']['total_applies']:,}")

        with col4:
            st.metric("Apply/Click %", f"{quartiles['middle_50']['apply_click_ratio']:.2f}%")

        with col5:
            st.metric("Avg Clicks/Vac", f"{quartiles['middle_50']['clicks_per_vacancy']:.1f}")

        with col6:
            st.metric("Avg Applies/Vac", f"{quartiles['middle_50']['applies_per_vacancy']:.2f}")

        st.markdown("---")

        # Row 4: Bottom 25% (Needs Improvement)
        st.markdown("### 🔴 Bottom 25% Performers")
        col1, col2, col3, col4, col5, col6 = st.columns(6)

        with col1:
            st.metric("Vacancies", f"{quartiles['bottom_25']['num_vacancies']:,}")

        with col2:
            st.metric("Total Clicks", f"{quartiles['bottom_25']['total_clicks']:,}")

        with col3:
            st.metric("Total Applies", f"{quartiles['bottom_25']['total_applies']:,}")

        with col4:
            st.metric("Apply/Click %", f"{quartiles['bottom_25']['apply_click_ratio']:.2f}%")

        with col5:
            st.metric("Avg Clicks/Vac", f"{quartiles['bottom_25']['clicks_per_vacancy']:.1f}")

        with col6:
            st.metric("Avg Applies/Vac", f"{quartiles['bottom_25']['applies_per_vacancy']:.2f}")

    else:
        # Fallback to simple view if not enough data
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric("Total Vacancies", f"{metrics['num_vacancies']:,}")

        with col2:
            st.metric("Total Clicks", f"{metrics['total_clicks']:,}")

        with col3:
            st.metric("Total Applies", f"{metrics['total_applies']:,}")

        with col4:
            st.metric("Apply/Click Ratio", f"{metrics['apply_click_ratio']:.2f}%")

    st.markdown("---")

    # Time Series (from pre-aggregated daily totals)
    if daily_totals is not None and len(daily_totals) > 0:
        st.subheader("Trends Over Time")

        # Check if any filters are active (beyond default)
        has_active_filters = False
        if 'overview_filters' in st.session_state and st.session_state.overview_filters:
            f = st.session_state.overview_filters
            if any([f.get('importer'), f.get('company'), f.get('region'),
                    f.get('occupation'), f.get('upgrades'), f.get('job_title')]):
                has_active_filters = True

        if has_active_filters:
            st.caption("Note: Trend data shows global site performance (not affected by filters above).")

        daily_data = daily_totals.copy()
        daily_data = daily_data.sort_values('event_date')

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=daily_data['event_date'], y=daily_data['clicks'],
                                name='Clicks', line=dict(color='blue')))
        fig.add_trace(go.Scatter(x=daily_data['event_date'], y=daily_data['applies'],
                                name='Applies', line=dict(color='green')))
        fig.update_layout(height=400, hovermode='x unified')
        st.plotly_chart(fig, width='stretch')

    st.markdown("---")

    # Performance by Dimension
    col1, col2 = st.columns(2)

    with col1:
        if 'importer_name' in filtered_df.columns:
            st.subheader("Performance by Importer")
            importer_stats = []
            for importer in filtered_df['importer_name'].unique():
                imp_df = filtered_df[filtered_df['importer_name'] == importer]
                imp_metrics = calculate_metrics(imp_df)
                importer_stats.append({
                    'Importer': importer,
                    'Vacancies': imp_metrics['num_vacancies'],
                    'Median Clicks': round(imp_metrics['median_clicks_per_vacancy'], 1),
                    'Mean Clicks': round(imp_metrics['mean_clicks_per_vacancy'], 1),
                    'Median Applies': round(imp_metrics['median_applies_per_vacancy'], 2),
                    'Mean Applies': round(imp_metrics['mean_applies_per_vacancy'], 2),
                    'Apply/Click %': round(imp_metrics['apply_click_ratio'], 2)
                })

            if importer_stats:
                importer_df = pd.DataFrame(importer_stats).sort_values('Mean Clicks', ascending=False)

                # Create grouped bar chart
                fig = go.Figure()
                fig.add_trace(go.Bar(
                    name='Median Clicks/Vacancy',
                    x=importer_df['Importer'],
                    y=importer_df['Median Clicks'],
                    text=importer_df['Median Clicks'],
                    textposition='auto',
                ))
                fig.add_trace(go.Bar(
                    name='Mean Clicks/Vacancy',
                    x=importer_df['Importer'],
                    y=importer_df['Mean Clicks'],
                    text=importer_df['Mean Clicks'],
                    textposition='auto',
                ))

                fig.update_layout(
                    barmode='group',
                    height=400,
                    xaxis_title='Importer',
                    yaxis_title='Clicks per Vacancy',
                    hovermode='x unified'
                )
                st.plotly_chart(fig, width='stretch')
            else:
                st.info("No importer data available for the selected filters.")

    with col2:
        _rdf = st.session_state.get('_region_df')
        has_region = (_rdf is not None and 'uk_region' in _rdf.columns) or 'uk_regions' in filtered_df.columns
        if has_region:
            st.subheader("Performance by Region")
            st.caption("A vacancy in multiple regions counts once per region.")
            region_stats = []

            if _rdf is not None and 'uk_region' in _rdf.columns:
                # Use pre-exploded region table
                from data.filters import apply_filters_to_region_data
                filtered_region = apply_filters_to_region_data(_rdf, st.session_state.get('overview_filters'))
                for region in sorted(filtered_region['uk_region'].dropna().unique()):
                    reg_df = filtered_region[filtered_region['uk_region'] == region]
                    reg_metrics = calculate_metrics(reg_df)
                    region_stats.append({
                        'Region': region,
                        'Vacancies': reg_metrics['num_vacancies'],
                        'Median Clicks': round(reg_metrics['median_clicks_per_vacancy'], 1),
                        'Mean Clicks': round(reg_metrics['mean_clicks_per_vacancy'], 1),
                        'Median Applies': round(reg_metrics['median_applies_per_vacancy'], 2),
                        'Mean Applies': round(reg_metrics['mean_applies_per_vacancy'], 2),
                        'Apply/Click %': round(reg_metrics['apply_click_ratio'], 2)
                    })
            else:
                # Fallback: pipe-split
                all_regions = set()
                for regions_str in filtered_df['uk_regions'].dropna():
                    for r in str(regions_str).split(' | '):
                        r = r.strip()
                        if r:
                            all_regions.add(r)
                for region in all_regions:
                    reg_mask = filtered_df['uk_regions'].apply(
                        lambda x: region in [r.strip() for r in str(x).split(' | ')] if pd.notna(x) else False
                    )
                    reg_df = filtered_df[reg_mask]
                    reg_metrics = calculate_metrics(reg_df)
                    region_stats.append({
                        'Region': region,
                        'Vacancies': reg_metrics['num_vacancies'],
                        'Median Clicks': round(reg_metrics['median_clicks_per_vacancy'], 1),
                        'Mean Clicks': round(reg_metrics['mean_clicks_per_vacancy'], 1),
                        'Median Applies': round(reg_metrics['median_applies_per_vacancy'], 2),
                        'Mean Applies': round(reg_metrics['mean_applies_per_vacancy'], 2),
                        'Apply/Click %': round(reg_metrics['apply_click_ratio'], 2)
                    })

            if region_stats:
                region_df = pd.DataFrame(region_stats).sort_values('Mean Clicks', ascending=False)

                # Create grouped bar chart
                fig = go.Figure()
                fig.add_trace(go.Bar(
                    name='Median Clicks/Vacancy',
                    x=region_df['Region'],
                    y=region_df['Median Clicks'],
                    text=region_df['Median Clicks'],
                    textposition='auto',
                ))
                fig.add_trace(go.Bar(
                    name='Mean Clicks/Vacancy',
                    x=region_df['Region'],
                    y=region_df['Mean Clicks'],
                    text=region_df['Mean Clicks'],
                    textposition='auto',
                ))

                fig.update_layout(
                    barmode='group',
                    height=400,
                    xaxis_title='Region',
                    yaxis_title='Clicks per Vacancy',
                    hovermode='x unified'
                )
                st.plotly_chart(fig, width='stretch')
            else:
                st.info("No region data available for the selected filters.")

    st.markdown("---")

    # Conversion Funnel
    st.subheader("Conversion Funnel")
    funnel_data = pd.DataFrame({
        'Stage': ['Vacancies', 'Clicks', 'Applies'],
        'Count': [metrics['num_vacancies'], metrics['total_clicks'], metrics['total_applies']]
    })
    fig = px.funnel(funnel_data, x='Count', y='Stage')
    fig.update_layout(height=400)
    st.plotly_chart(fig, width='stretch')

# ============================================================================
# TAB 2: DEEP DIVE
# ============================================================================

def create_deep_dive_tab(df):
    """Create the Deep Dive tab."""
    st.header("🔍 Deep Dive")

    # Filters
    with st.expander("🔍 Filters", expanded=True):
        filters, apply_clicked = create_filter_panel(df, 'deepdive')

    # Apply filters
    if apply_clicked or 'deepdive_filters' in st.session_state:
        if apply_clicked:
            st.session_state.deepdive_filters = filters
        filtered_df = apply_filters_to_data(df, st.session_state.deepdive_filters)
    else:
        filtered_df = df.copy()

    # Benchmark Comparison Table
    st.subheader("📊 Benchmark Comparison Table")

    dimension = st.selectbox(
        "Group by:",
        ['Importer', 'Region', 'Occupation', 'Company'],
        key='deepdive_dimension'
    )

    column_map = {
        'Importer': 'importer_name',
        'Region': 'uk_regions',
        'Occupation': 'occupation',
        'Company': 'organization_name'
    }

    col_name = column_map[dimension]
    if col_name in filtered_df.columns:
        benchmark_data = []

        if dimension == 'Region':
            _rdf = st.session_state.get('_region_df')
            if _rdf is not None and 'uk_region' in _rdf.columns:
                from data.filters import apply_filters_to_region_data
                filtered_region = apply_filters_to_region_data(_rdf, st.session_state.get('deepdive_filters'))
                for value in sorted(filtered_region['uk_region'].dropna().unique()):
                    subset = filtered_region[filtered_region['uk_region'] == value]
                    metrics = calculate_metrics(subset)
                    benchmark_data.append({
                        dimension: value,
                        'Vacancies': metrics['num_vacancies'],
                        'Total Clicks': metrics['total_clicks'],
                        'Total Applies': metrics['total_applies'],
                        'Apply/Click %': round(metrics['apply_click_ratio'], 2),
                        'Median Clicks/Vac': round(metrics['median_clicks_per_vacancy'], 1),
                        'Mean Clicks/Vac': round(metrics['mean_clicks_per_vacancy'], 1),
                        'Median Applies/Vac': round(metrics['median_applies_per_vacancy'], 2),
                        'Mean Applies/Vac': round(metrics['mean_applies_per_vacancy'], 2)
                    })
            else:
                # Fallback: pipe-split
                all_values = set()
                for regions_str in filtered_df[col_name].dropna():
                    for r in str(regions_str).split(' | '):
                        r = r.strip()
                        if r:
                            all_values.add(r)
                for value in all_values:
                    mask = filtered_df[col_name].apply(
                        lambda x: value in [r.strip() for r in str(x).split(' | ')] if pd.notna(x) else False
                    )
                    subset = filtered_df[mask]
                    metrics = calculate_metrics(subset)
                    benchmark_data.append({
                        dimension: value,
                        'Vacancies': metrics['num_vacancies'],
                        'Total Clicks': metrics['total_clicks'],
                        'Total Applies': metrics['total_applies'],
                        'Apply/Click %': round(metrics['apply_click_ratio'], 2),
                        'Median Clicks/Vac': round(metrics['median_clicks_per_vacancy'], 1),
                        'Mean Clicks/Vac': round(metrics['mean_clicks_per_vacancy'], 1),
                        'Median Applies/Vac': round(metrics['median_applies_per_vacancy'], 2),
                        'Mean Applies/Vac': round(metrics['mean_applies_per_vacancy'], 2)
                    })
        else:
            for value in filtered_df[col_name].unique():
                subset = filtered_df[filtered_df[col_name] == value]
                metrics = calculate_metrics(subset)
                benchmark_data.append({
                    dimension: value,
                    'Vacancies': metrics['num_vacancies'],
                    'Total Clicks': metrics['total_clicks'],
                    'Total Applies': metrics['total_applies'],
                    'Apply/Click %': round(metrics['apply_click_ratio'], 2),
                    'Median Clicks/Vac': round(metrics['median_clicks_per_vacancy'], 1),
                    'Mean Clicks/Vac': round(metrics['mean_clicks_per_vacancy'], 1),
                    'Median Applies/Vac': round(metrics['median_applies_per_vacancy'], 2),
                    'Mean Applies/Vac': round(metrics['mean_applies_per_vacancy'], 2)
                })

        if benchmark_data:
            benchmark_df = pd.DataFrame(benchmark_data).sort_values('Mean Clicks/Vac', ascending=False)
            st.dataframe(benchmark_df, width='stretch')
        else:
            st.info("No benchmark data available for the selected filters.")
            benchmark_df = pd.DataFrame()

        # Export
        csv = benchmark_df.to_csv(index=False).encode('utf-8') if not benchmark_df.empty else b''
        st.download_button(
            "📥 Download Benchmark Data",
            csv,
            f"benchmark_{dimension.lower()}_{datetime.now().strftime('%Y%m%d')}.csv",
            "text/csv"
        )

    st.markdown("---")

    # Heatmap
    st.subheader("🗺️ Performance Heatmap")

    _rdf = st.session_state.get('_region_df')
    has_heatmap = (
        (_rdf is not None and 'uk_region' in _rdf.columns)
        or 'uk_regions' in filtered_df.columns
    ) and 'importer_name' in filtered_df.columns
    if has_heatmap:
        heatmap_data = []

        if _rdf is not None and 'uk_region' in _rdf.columns:
            from data.filters import apply_filters_to_region_data
            filtered_region = apply_filters_to_region_data(_rdf, st.session_state.get('deepdive_filters'))
            for region in sorted(filtered_region['uk_region'].dropna().unique()):
                reg_subset = filtered_region[filtered_region['uk_region'] == region]
                for importer in reg_subset['importer_name'].unique():
                    subset = reg_subset[reg_subset['importer_name'] == importer]
                    if len(subset) > 0:
                        metrics = calculate_metrics(subset)
                        heatmap_data.append({
                            'Region': region,
                            'Importer': importer,
                            'Clicks/Vacancy': metrics['clicks_per_vacancy'],
                            'Applies/Vacancy': metrics['applies_per_vacancy'],
                            'Apply/Click %': metrics['apply_click_ratio']
                        })
        else:
            all_regions = set()
            for regions_str in filtered_df['uk_regions'].dropna():
                for r in str(regions_str).split(' | '):
                    r = r.strip()
                    if r:
                        all_regions.add(r)
            for region in all_regions:
                reg_mask = filtered_df['uk_regions'].apply(
                    lambda x: region in [r.strip() for r in str(x).split(' | ')] if pd.notna(x) else False
                )
                for importer in filtered_df['importer_name'].unique():
                    subset = filtered_df[
                        reg_mask &
                        (filtered_df['importer_name'] == importer)
                    ]
                    if len(subset) > 0:
                        metrics = calculate_metrics(subset)
                        heatmap_data.append({
                            'Region': region,
                            'Importer': importer,
                            'Clicks/Vacancy': metrics['clicks_per_vacancy'],
                            'Applies/Vacancy': metrics['applies_per_vacancy'],
                            'Apply/Click %': metrics['apply_click_ratio']
                        })

        if heatmap_data:
            heatmap_df = pd.DataFrame(heatmap_data)

            # Allow user to select metric for heatmap
            heatmap_metric = st.selectbox(
                "Select metric for heatmap:",
                ['Clicks/Vacancy', 'Applies/Vacancy', 'Apply/Click %'],
                key='heatmap_metric'
            )

            heatmap_pivot = heatmap_df.pivot(index='Region', columns='Importer', values=heatmap_metric)

            fig = px.imshow(
                heatmap_pivot,
                labels=dict(x="Importer", y="Region", color=heatmap_metric),
                aspect="auto",
                color_continuous_scale='RdYlGn'
            )
            fig.update_layout(height=500)
            st.plotly_chart(fig, width='stretch')

# ============================================================================
# TAB 3: VACANCY PERFORMANCE
# ============================================================================

def create_vacancy_performance_tab(df, full_df=None):
    """Create the Vacancy Performance tab.

    Args:
        df: The main dataframe (used for filtering)
        full_df: The full unfiltered dataframe (used for benchmark calculations)
    """
    st.header("📋 Vacancy Performance")

    # Use full_df if provided, otherwise use df
    if full_df is None:
        full_df = df

    # Filters
    with st.expander("🔍 Filters", expanded=True):
        filters, apply_clicked = create_filter_panel(df, 'vacancy')

    # Apply filters
    if apply_clicked or 'vacancy_filters' in st.session_state:
        if apply_clicked:
            st.session_state.vacancy_filters = filters
        filtered_df = apply_filters_to_data(df, st.session_state.vacancy_filters)
    else:
        filtered_df = df.copy()

    # Each row in filtered_df is already a vacancy with clicks/applies columns
    job_col = 'entity_id' if 'entity_id' in filtered_df.columns else filtered_df.columns[0]

    vacancy_data = []
    for _, job in filtered_df.iterrows():
        job_id = job[job_col]
        clicks = int(job.get('clicks', 0))
        applies = int(job.get('applies', 0))
        ratio = (applies / clicks * 100) if clicks > 0 else 0

        # Get vacancy status from workflow_state
        status = job.get('workflow_state', 'Unknown')
        is_published = status == 'published'

        # Calculate days active with improved logic
        days_active = None
        start_date = job.get('start_date')
        end_date = job.get('end_date')

        if pd.notna(start_date):
            if pd.notna(end_date):
                days_active = (end_date - start_date).days
            elif is_published:
                today = pd.Timestamp(datetime.now())
                days_active = (today - start_date).days

        # Get occupation from occupational_fields
        occupation = job.get('occupational_fields', 'Unknown')
        if pd.notna(occupation) and str(occupation).strip():
            occupation = str(occupation).split('|')[0].strip()
        else:
            occupation = 'Unknown'

        # Get upgrades
        upgrades_str = ', '.join(job.get('upgrades_list', [])) if 'upgrades_list' in job.index else ''

        vacancy_data.append({
            'Title': job.get('title', job.get('organization_name', 'Unknown')),
            'Company': job.get('organization_name', 'Unknown'),
            'Job ID': job_id,
            'Status': status,
            'Start Date': start_date if pd.notna(start_date) else None,
            'End Date': end_date if pd.notna(end_date) else None,
            'Days Active': int(days_active) if days_active is not None and days_active > 0 else None,
            'Region': job.get('uk_regions', 'Unknown'),
            'Occupation': occupation,
            'Importer': job.get('importer_name', 'Unknown'),
            'Upgrades': upgrades_str if upgrades_str else 'None',
            'Clicks': clicks,
            'Applies': applies,
            'Ratio %': round(ratio, 2) if clicks > 0 else None,
            'Clicks/Day': round(clicks / days_active, 2) if days_active and days_active > 0 else None,
            'Applies/Day': round(applies / days_active, 2) if days_active and days_active > 0 else None
        })

    vacancy_df = pd.DataFrame(vacancy_data)

    # Check if we have any data
    if len(vacancy_df) == 0:
        st.warning("⚠️ No vacancy data found for the selected filters. Try adjusting your date range or filters.")
        return

    # Calculate occupation benchmarks from FULL dataset (static benchmarks)
    # With pre-aggregated data, each row already has clicks/applies
    full_job_col = 'entity_id' if 'entity_id' in full_df.columns else full_df.columns[0]

    # Build occupation stats directly from the full dataframe
    full_df_with_occ = full_df.copy()
    full_df_with_occ['_occupation'] = full_df_with_occ['occupational_fields'].apply(
        lambda x: str(x).split('|')[0].strip() if pd.notna(x) and str(x).strip() else 'Unknown'
    )

    occupation_stats = {}
    for occupation in full_df_with_occ['_occupation'].unique():
        occ_vacancies = full_df_with_occ[full_df_with_occ['_occupation'] == occupation]
        occupation_stats[occupation] = {
            'avg_clicks': occ_vacancies['clicks'].mean() if 'clicks' in occ_vacancies.columns else 0,
            'avg_applies': occ_vacancies['applies'].mean() if 'applies' in occ_vacancies.columns else 0
        }

    # Add occupation benchmarks to filtered vacancy data
    if 'Occupation' in vacancy_df.columns:
        vacancy_df['Avg Clicks (Occupation)'] = vacancy_df['Occupation'].map(
            lambda x: round(occupation_stats.get(x, {}).get('avg_clicks', 0), 1)
        )
        vacancy_df['Avg Applies (Occupation)'] = vacancy_df['Occupation'].map(
            lambda x: round(occupation_stats.get(x, {}).get('avg_applies', 0), 1)
        )

    vacancy_df = vacancy_df.sort_values('Clicks', ascending=False)

    # Metrics
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Total Vacancies", len(vacancy_df))
    with col2:
        st.metric("Total Clicks", f"{vacancy_df['Clicks'].sum():,}")
    with col3:
        st.metric("Total Applies", f"{vacancy_df['Applies'].sum():,}")
    with col4:
        avg_ratio = vacancy_df['Ratio %'].mean()
        st.metric("Avg Apply Rate", f"{avg_ratio:.2f}%")

    # Display table
    st.subheader(f"Vacancy Data ({len(vacancy_df)} vacancies)")
    st.dataframe(vacancy_df, width='stretch', height=600, hide_index=True)

    # Export
    csv = vacancy_df.to_csv(index=False).encode('utf-8')
    st.download_button(
        "📥 Download Vacancy Report",
        csv,
        f"vacancy_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
        "text/csv"
    )

# ============================================================================
# TAB 4: COMPARISON
# ============================================================================

def create_comparison_tab(df):
    """Create the Comparison tab."""
    st.header("⚖️ Comparison")
    st.info("💡 Select filters for each side, then click Apply Filters to compare")

    col_left, col_right = st.columns(2)

    # Left Side
    with col_left:
        st.subheader("📊 Side A")
        with st.expander("🔍 Filters", expanded=True):
            filters_left, apply_left = create_filter_panel(df, 'comp_left')

        if apply_left or 'comp_left_filters' in st.session_state:
            if apply_left:
                st.session_state.comp_left_filters = filters_left
            filtered_left = apply_filters_to_data(df, st.session_state.comp_left_filters)
        else:
            filtered_left = df.copy()

        metrics_left = calculate_metrics(filtered_left)

        st.markdown("### Totals")
        st.metric("Vacancies", f"{metrics_left['num_vacancies']:,}")
        st.metric("Clicks", f"{metrics_left['total_clicks']:,}")
        st.metric("Applies", f"{metrics_left['total_applies']:,}")
        st.metric("Apply/Click %", f"{metrics_left['apply_click_ratio']:.2f}%")
        st.metric("Clicks/Vacancy", f"{metrics_left['clicks_per_vacancy']:.1f}")
        st.metric("Applies/Vacancy", f"{metrics_left['applies_per_vacancy']:.2f}")

    # Right Side
    with col_right:
        st.subheader("📊 Side B")
        with st.expander("🔍 Filters", expanded=True):
            filters_right, apply_right = create_filter_panel(df, 'comp_right')

        if apply_right or 'comp_right_filters' in st.session_state:
            if apply_right:
                st.session_state.comp_right_filters = filters_right
            filtered_right = apply_filters_to_data(df, st.session_state.comp_right_filters)
        else:
            filtered_right = df.copy()

        metrics_right = calculate_metrics(filtered_right)

        st.markdown("### Totals")
        st.metric("Vacancies", f"{metrics_right['num_vacancies']:,}")
        st.metric("Clicks", f"{metrics_right['total_clicks']:,}")
        st.metric("Applies", f"{metrics_right['total_applies']:,}")
        st.metric("Apply/Click %", f"{metrics_right['apply_click_ratio']:.2f}%")
        st.metric("Clicks/Vacancy", f"{metrics_right['clicks_per_vacancy']:.1f}")
        st.metric("Applies/Vacancy", f"{metrics_right['applies_per_vacancy']:.2f}")

    # Comparison Summary
    st.markdown("---")
    st.subheader("📈 Comparison Summary")

    col1, col2, col3, col4, col5, col6 = st.columns(6)

    with col1:
        st.markdown("**Vacancies**")
        diff = metrics_right['num_vacancies'] - metrics_left['num_vacancies']
        pct = ((metrics_right['num_vacancies'] / metrics_left['num_vacancies'] - 1) * 100) if metrics_left['num_vacancies'] > 0 else 0
        st.markdown(f"Side A: {metrics_left['num_vacancies']:,}")
        st.markdown(f"Side B: {metrics_right['num_vacancies']:,}")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+,} ({pct:+.1f}%)")

    with col2:
        st.markdown("**Clicks**")
        diff = metrics_right['total_clicks'] - metrics_left['total_clicks']
        pct = ((metrics_right['total_clicks'] / metrics_left['total_clicks'] - 1) * 100) if metrics_left['total_clicks'] > 0 else 0
        st.markdown(f"Side A: {metrics_left['total_clicks']:,}")
        st.markdown(f"Side B: {metrics_right['total_clicks']:,}")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+,} ({pct:+.1f}%)")

    with col3:
        st.markdown("**Applies**")
        diff = metrics_right['total_applies'] - metrics_left['total_applies']
        pct = ((metrics_right['total_applies'] / metrics_left['total_applies'] - 1) * 100) if metrics_left['total_applies'] > 0 else 0
        st.markdown(f"Side A: {metrics_left['total_applies']:,}")
        st.markdown(f"Side B: {metrics_right['total_applies']:,}")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+,} ({pct:+.1f}%)")

    with col4:
        st.markdown("**Apply/Click %**")
        diff = metrics_right['apply_click_ratio'] - metrics_left['apply_click_ratio']
        st.markdown(f"Side A: {metrics_left['apply_click_ratio']:.2f}%")
        st.markdown(f"Side B: {metrics_right['apply_click_ratio']:.2f}%")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+.2f}%")

    with col5:
        st.markdown("**Clicks/Vacancy**")
        diff = metrics_right['clicks_per_vacancy'] - metrics_left['clicks_per_vacancy']
        pct = ((metrics_right['clicks_per_vacancy'] / metrics_left['clicks_per_vacancy'] - 1) * 100) if metrics_left['clicks_per_vacancy'] > 0 else 0
        st.markdown(f"Side A: {metrics_left['clicks_per_vacancy']:.1f}")
        st.markdown(f"Side B: {metrics_right['clicks_per_vacancy']:.1f}")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+.1f} ({pct:+.1f}%)")

    with col6:
        st.markdown("**Applies/Vacancy**")
        diff = metrics_right['applies_per_vacancy'] - metrics_left['applies_per_vacancy']
        pct = ((metrics_right['applies_per_vacancy'] / metrics_left['applies_per_vacancy'] - 1) * 100) if metrics_left['applies_per_vacancy'] > 0 else 0
        st.markdown(f"Side A: {metrics_left['applies_per_vacancy']:.2f}")
        st.markdown(f"Side B: {metrics_right['applies_per_vacancy']:.2f}")
        color = "🟢" if diff > 0 else "🔴" if diff < 0 else "⚪"
        st.markdown(f"{color} Diff: {diff:+.2f} ({pct:+.1f}%)")

    # Visual comparison
    st.markdown("---")
    st.subheader("Visual Comparison")

    comparison_data = pd.DataFrame({
        'Side': ['A', 'B', 'A', 'B', 'A', 'B'],
        'Metric': ['Vacancies', 'Vacancies', 'Clicks', 'Clicks', 'Applies', 'Applies'],
        'Value': [
            metrics_left['num_vacancies'], metrics_right['num_vacancies'],
            metrics_left['total_clicks'], metrics_right['total_clicks'],
            metrics_left['total_applies'], metrics_right['total_applies']
        ]
    })

    fig = px.bar(comparison_data, x='Metric', y='Value', color='Side', barmode='group')
    fig.update_layout(height=400)
    st.plotly_chart(fig, width='stretch')

# ============================================================================
# TAB 5: SALES INTELLIGENCE
# ============================================================================

def create_sales_intelligence_tab(df):
    """Create the Sales Intelligence tab with upgrade ROI, client scorecards,
    underperforming alerts, and sector benchmarks."""

    with st.expander("🔍 Filters", expanded=False):
        filters, apply_clicked = create_filter_panel(df, 'sales')

    if apply_clicked:
        st.session_state['sales_filters'] = filters

    filtered_df = apply_filters_to_data(df, st.session_state.get('sales_filters'))

    # ------------------------------------------------------------------
    # Section 1: Upgrade ROI
    # ------------------------------------------------------------------
    st.subheader("1. Upgrade ROI by Occupation")

    st.info(
        "**Methodology:** Figures show average clicks/applies per vacancy for each occupation and upgrade type. "
        "Only vacancies with a **single upgrade** (or no upgrade) are included to avoid double-counting. "
        "Vacancy counts shown in brackets — figures from small samples (n<10) are marked with * and should be treated with caution. "
        "Occupations with fewer than 5 vacancies are excluded."
    )

    if len(filtered_df) > 0 and 'upgrades_list' in filtered_df.columns:
        # Classify each vacancy
        work_df = filtered_df.copy()
        work_df['upgrade_count'] = work_df['upgrades_list'].apply(len)
        work_df['upgrade_category'] = work_df['upgrades_list'].apply(
            lambda x: 'No Upgrade' if not x else (x[0] if len(x) == 1 else ' + '.join(sorted(x)))
        )

        # ---- Primary view: single-upgrade vacancies only ----
        single_df = work_df[work_df['upgrade_count'] <= 1].copy()

        if len(single_df) > 0:
            # Aggregate: occupation × upgrade → mean clicks/applies + count
            agg = single_df.groupby(['occupation', 'upgrade_category']).agg(
                avg_clicks=('clicks', 'mean'),
                avg_applies=('applies', 'mean'),
                n=('clicks', 'count')
            ).reset_index()

            # Pivot to tables
            clicks_table = agg.pivot(index='occupation', columns='upgrade_category', values='avg_clicks').round(1)
            applies_table = agg.pivot(index='occupation', columns='upgrade_category', values='avg_applies').round(1)
            count_table = agg.pivot(index='occupation', columns='upgrade_category', values='n').fillna(0).astype(int)

            # Ensure "No Upgrade" is first column
            for tbl in [clicks_table, applies_table, count_table]:
                cols = list(tbl.columns)
                if 'No Upgrade' in cols:
                    cols.remove('No Upgrade')
                    cols = ['No Upgrade'] + sorted(cols)

            clicks_table = clicks_table.reindex(columns=cols)
            applies_table = applies_table.reindex(columns=cols)
            count_table = count_table.reindex(columns=cols).fillna(0).astype(int)

            # Filter: 5+ vacancies per occupation
            total_per_occ = count_table.sum(axis=1)
            mask = total_per_occ >= 5
            clicks_table = clicks_table[mask]
            applies_table = applies_table[mask]
            count_table = count_table[mask]

            if len(clicks_table) > 0:
                # Sort by total vacancy count descending
                sort_order = count_table.sum(axis=1).sort_values(ascending=False).index
                clicks_table = clicks_table.loc[sort_order]
                applies_table = applies_table.loc[sort_order]
                count_table = count_table.loc[sort_order]

                # Build display tables with "value (n=X)" format
                def build_display_table(val_table, cnt_table):
                    display = val_table.copy().astype(str)
                    for col in display.columns:
                        for idx in display.index:
                            v = val_table.at[idx, col]
                            n = int(cnt_table.at[idx, col]) if idx in cnt_table.index and col in cnt_table.columns else 0
                            if pd.isna(v) or n == 0:
                                display.at[idx, col] = '—'
                            else:
                                flag = '*' if n < 10 else ''
                                display.at[idx, col] = f'{v:.1f} (n={n}){flag}'
                    return display

                # Build uplift tables (% vs No Upgrade baseline)
                def build_uplift_table(val_table):
                    uplift = val_table.copy()
                    if 'No Upgrade' not in uplift.columns:
                        return None
                    for col in uplift.columns:
                        if col == 'No Upgrade':
                            uplift[col] = '—'
                        else:
                            for idx in uplift.index:
                                v = val_table.at[idx, col]
                                baseline = val_table.at[idx, 'No Upgrade']
                                if pd.isna(v) or pd.isna(baseline) or baseline == 0:
                                    uplift.at[idx, col] = '—'
                                else:
                                    pct = ((v - baseline) / baseline) * 100
                                    uplift.at[idx, col] = f'{pct:+.0f}%'
                    return uplift

                # View toggle
                view_mode = st.radio(
                    "Display mode", ["Absolute values", "% uplift vs No Upgrade"],
                    horizontal=True, key='upgrade_roi_mode'
                )

                tab_clicks, tab_applies = st.tabs(["Clicks per Vacancy", "Applies per Vacancy"])

                def style_display(display_df, val_table):
                    """Apply green/red highlighting vs No Upgrade baseline."""
                    def highlight_row(row):
                        baseline_key = 'No Upgrade'
                        styles = []
                        for col in row.index:
                            cell = row[col]
                            if col == baseline_key:
                                styles.append('background-color: #f0f0f0')
                            elif cell == '—' or baseline_key not in val_table.columns:
                                styles.append('')
                            else:
                                idx = row.name
                                v = val_table.at[idx, col] if col in val_table.columns else None
                                b = val_table.at[idx, baseline_key] if baseline_key in val_table.columns else None
                                if pd.isna(v) or pd.isna(b) or b == 0:
                                    styles.append('')
                                elif v > b * 1.1:
                                    styles.append('background-color: #c6efce')
                                elif v < b * 0.9:
                                    styles.append('background-color: #ffc7ce')
                                else:
                                    styles.append('')
                        return styles
                    return display_df.style.apply(highlight_row, axis=1)

                with tab_clicks:
                    if view_mode == "Absolute values":
                        display = build_display_table(clicks_table, count_table)
                    else:
                        display = build_uplift_table(clicks_table)
                    if display is not None:
                        styled = style_display(display, clicks_table)
                        st.dataframe(styled, width='stretch', height=min(500, 35 * len(display) + 40))

                with tab_applies:
                    if view_mode == "Absolute values":
                        display = build_display_table(applies_table, count_table)
                    else:
                        display = build_uplift_table(applies_table)
                    if display is not None:
                        styled = style_display(display, applies_table)
                        st.dataframe(styled, width='stretch', height=min(500, 35 * len(display) + 40))

                st.caption("Green = >10% above 'No Upgrade' baseline | Red = >10% below | Grey = baseline | * = small sample (n<10)")

                # ---- Secondary view: multi-upgrade combinations ----
                multi_df = work_df[work_df['upgrade_count'] > 1]
                if len(multi_df) > 0:
                    with st.expander(f"Upgrade Combinations ({len(multi_df)} vacancies with multiple upgrades)"):
                        combo_agg = multi_df.groupby('upgrade_category').agg(
                            vacancies=('clicks', 'count'),
                            avg_clicks=('clicks', 'mean'),
                            avg_applies=('applies', 'mean')
                        ).reset_index().sort_values('vacancies', ascending=False)
                        combo_agg.columns = ['Upgrade Combination', 'Vacancies', 'Avg Clicks', 'Avg Applies']

                        # Add baseline comparison
                        no_upgrade_clicks = single_df[single_df['upgrade_category'] == 'No Upgrade']['clicks'].mean()
                        no_upgrade_applies = single_df[single_df['upgrade_category'] == 'No Upgrade']['applies'].mean()
                        if pd.notna(no_upgrade_clicks) and no_upgrade_clicks > 0:
                            combo_agg['Clicks vs No Upgrade'] = combo_agg['Avg Clicks'].apply(
                                lambda x: f'{((x / no_upgrade_clicks) - 1) * 100:+.0f}%' if pd.notna(x) else '—')
                        if pd.notna(no_upgrade_applies) and no_upgrade_applies > 0:
                            combo_agg['Applies vs No Upgrade'] = combo_agg['Avg Applies'].apply(
                                lambda x: f'{((x / no_upgrade_applies) - 1) * 100:+.0f}%' if pd.notna(x) else '—')

                        st.dataframe(
                            combo_agg.style.format({'Avg Clicks': '{:.1f}', 'Avg Applies': '{:.1f}'}),
                            width='stretch', hide_index=True
                        )
                        st.caption("These vacancies had multiple upgrades applied simultaneously. "
                                   "Performance shown is for the full combination, not individual upgrade types.")
            else:
                st.info("Not enough data (need 5+ vacancies per occupation).")
    else:
        st.info("No data available for upgrade analysis.")

    st.markdown("---")

    # ------------------------------------------------------------------
    # Section 2: Client Scorecards
    # ------------------------------------------------------------------
    st.subheader("2. Client Scorecards")
    st.caption("Per-organisation performance benchmarked against their sector average.")

    if len(filtered_df) > 0 and 'organization_name' in filtered_df.columns:
        # Organisation selector
        orgs = sorted(filtered_df['organization_name'].dropna().unique())
        selected_orgs = st.multiselect("Select organisations to view", orgs, key='sales_org_select')

        # Calculate sector (occupation) benchmarks
        sector_benchmarks = filtered_df.groupby('occupation').agg(
            sector_avg_clicks=('clicks', 'mean'),
            sector_avg_applies=('applies', 'mean'),
            sector_vacancy_count=('clicks', 'count')
        ).reset_index()

        # Calculate per-org metrics
        org_data = filtered_df.groupby(['organization_name', 'occupation']).agg(
            vacancies=('clicks', 'count'),
            total_clicks=('clicks', 'sum'),
            total_applies=('applies', 'sum'),
            avg_clicks=('clicks', 'mean'),
            avg_applies=('applies', 'mean')
        ).reset_index()

        org_data = org_data.merge(sector_benchmarks, on='occupation', how='left')
        org_data['apply_rate'] = (org_data['total_applies'] / org_data['total_clicks'] * 100).fillna(0).round(1)
        org_data['clicks_vs_sector'] = ((org_data['avg_clicks'] / org_data['sector_avg_clicks'] - 1) * 100).round(1)
        org_data['applies_vs_sector'] = ((org_data['avg_applies'] / org_data['sector_avg_applies'] - 1) * 100).round(1)

        display_orgs = selected_orgs if selected_orgs else orgs[:10]

        for org in display_orgs:
            org_rows = org_data[org_data['organization_name'] == org]
            if len(org_rows) == 0:
                continue

            total_vacancies = int(org_rows['vacancies'].sum())
            total_clicks = int(org_rows['total_clicks'].sum())
            total_applies = int(org_rows['total_applies'].sum())
            apply_rate = (total_applies / total_clicks * 100) if total_clicks > 0 else 0

            with st.expander(f"**{org}** — {total_vacancies} vacancies, {total_clicks:,} clicks, {total_applies:,} applies"):
                col1, col2, col3, col4 = st.columns(4)
                col1.metric("Vacancies", f"{total_vacancies:,}")
                col2.metric("Total Clicks", f"{total_clicks:,}")
                col3.metric("Total Applies", f"{total_applies:,}")
                col4.metric("Apply Rate", f"{apply_rate:.1f}%")

                # Per-occupation breakdown with sector comparison
                display_cols = ['occupation', 'vacancies', 'avg_clicks', 'sector_avg_clicks', 'clicks_vs_sector',
                                'avg_applies', 'sector_avg_applies', 'applies_vs_sector']
                display_df = org_rows[display_cols].copy()
                display_df.columns = ['Occupation', 'Vacancies', 'Avg Clicks', 'Sector Avg Clicks', 'Clicks vs Sector %',
                                      'Avg Applies', 'Sector Avg Applies', 'Applies vs Sector %']
                display_df = display_df.sort_values('Vacancies', ascending=False)

                def color_vs_sector(val):
                    if pd.isna(val):
                        return ''
                    if val > 10:
                        return 'color: #006100; background-color: #c6efce'
                    elif val < -10:
                        return 'color: #9c0006; background-color: #ffc7ce'
                    return ''

                styled = display_df.style.map(
                    color_vs_sector, subset=['Clicks vs Sector %', 'Applies vs Sector %']
                ).format({
                    'Avg Clicks': '{:.1f}', 'Sector Avg Clicks': '{:.1f}',
                    'Avg Applies': '{:.1f}', 'Sector Avg Applies': '{:.1f}',
                    'Clicks vs Sector %': '{:+.1f}%', 'Applies vs Sector %': '{:+.1f}%'
                })
                st.dataframe(styled, width='stretch', hide_index=True)
    else:
        st.info("No data available for client scorecards.")

    st.markdown("---")

    # ------------------------------------------------------------------
    # Section 3: Underperforming Vacancy Alerts
    # ------------------------------------------------------------------
    st.subheader("3. Underperforming Vacancy Alerts")
    st.caption("Published vacancies performing below the 25th percentile for their occupation.")

    if len(filtered_df) > 0:
        # Only look at published vacancies
        published = filtered_df[filtered_df['workflow_state'] == 'published'].copy()

        if len(published) > 0:
            # Calculate per-occupation thresholds
            occ_stats = filtered_df.groupby('occupation').agg(
                p25_clicks=('clicks', lambda x: x.quantile(0.25)),
                avg_clicks=('clicks', 'mean'),
                avg_applies=('applies', 'mean')
            ).reset_index()

            published = published.merge(occ_stats, on='occupation', how='left')

            # Flag underperformers
            underperformers = published[published['clicks'] < published['p25_clicks']].copy()

            if len(underperformers) > 0:
                # Calculate days live
                now = pd.Timestamp.now()
                if 'first_event_date' in underperformers.columns:
                    underperformers['days_live'] = (now - underperformers['first_event_date']).dt.days

                display_df = underperformers[['title', 'organization_name', 'occupation',
                                              'clicks', 'avg_clicks', 'applies', 'avg_applies',
                                              'days_live']].copy()
                display_df.columns = ['Title', 'Organisation', 'Occupation',
                                      'Clicks', 'Sector Avg Clicks', 'Applies', 'Sector Avg Applies',
                                      'Days Live']
                display_df = display_df.sort_values('Clicks', ascending=True)

                st.warning(f"**{len(underperformers):,}** published vacancies are below the 25th percentile for their occupation.")

                st.dataframe(
                    display_df.style.format({
                        'Sector Avg Clicks': '{:.1f}',
                        'Sector Avg Applies': '{:.1f}'
                    }),
                    width='stretch',
                    hide_index=True,
                    height=min(600, 35 * len(display_df) + 40)
                )
            else:
                st.success("No underperforming published vacancies found.")
        else:
            st.info("No published vacancies in the filtered data.")
    else:
        st.info("No data available.")

    st.markdown("---")

    # ------------------------------------------------------------------
    # Section 4: Sector Benchmarks
    # ------------------------------------------------------------------
    st.subheader("4. Sector Benchmarks")
    st.caption("Performance benchmarks by occupation — use these figures when talking to clients.")

    if len(filtered_df) > 0 and 'occupation' in filtered_df.columns:
        sector = filtered_df.groupby('occupation').agg(
            vacancies=('clicks', 'count'),
            median_clicks=('clicks', 'median'),
            mean_clicks=('clicks', 'mean'),
            median_applies=('applies', 'median'),
            mean_applies=('applies', 'mean'),
        ).reset_index()

        # Apply rate per occupation
        occ_totals = filtered_df.groupby('occupation').agg(
            total_clicks=('clicks', 'sum'),
            total_applies=('applies', 'sum')
        ).reset_index()
        occ_totals['apply_rate'] = (occ_totals['total_applies'] / occ_totals['total_clicks'] * 100).fillna(0)
        sector = sector.merge(occ_totals[['occupation', 'apply_rate']], on='occupation', how='left')

        # Average days live
        if 'first_event_date' in filtered_df.columns and 'last_event_date' in filtered_df.columns:
            filtered_df_copy = filtered_df.copy()
            filtered_df_copy['days_live'] = (filtered_df_copy['last_event_date'] - filtered_df_copy['first_event_date']).dt.days
            avg_days = filtered_df_copy.groupby('occupation')['days_live'].mean().reset_index()
            avg_days.columns = ['occupation', 'avg_days_live']
            sector = sector.merge(avg_days, on='occupation', how='left')
        else:
            sector['avg_days_live'] = None

        sector = sector[sector['vacancies'] >= 5].sort_values('vacancies', ascending=False)
        sector.columns = ['Occupation', 'Vacancies', 'Median Clicks', 'Mean Clicks',
                          'Median Applies', 'Mean Applies', 'Apply Rate %', 'Avg Days Live']

        if len(sector) > 0:
            def _gradient_green(val, col_min, col_max):
                if pd.isna(val) or col_max == col_min:
                    return ''
                pct = (val - col_min) / (col_max - col_min)
                r = int(245 - pct * 47)
                g = int(245 - pct * 12)
                b = int(245 - pct * 47)
                return f'background-color: rgb({r},{g},{b})'

            def apply_manual_gradient(df_styled, col, color='green'):
                col_data = sector[col].dropna()
                if len(col_data) == 0:
                    return df_styled
                cmin, cmax = col_data.min(), col_data.max()

                def _grad(v, r, g, b, cmin=cmin, cmax=cmax):
                    """Safe gradient that handles pd.NA without boolean ambiguity."""
                    try:
                        if pd.isna(v) or cmax <= cmin:
                            return ''
                        intensity = min((float(v) - float(cmin)) / (float(cmax) - float(cmin)), 1) * 0.3
                        return f'background-color: rgba({r},{g},{b},{intensity:.2f})'
                    except (TypeError, ValueError):
                        return ''

                if color == 'green':
                    return df_styled.map(lambda v: _grad(v, 34, 139, 34), subset=[col])
                elif color == 'blue':
                    return df_styled.map(lambda v: _grad(v, 30, 90, 200), subset=[col])
                else:
                    return df_styled.map(lambda v: _grad(v, 220, 120, 20), subset=[col])

            styled = sector.style.format({
                'Median Clicks': '{:.0f}', 'Mean Clicks': '{:.1f}',
                'Median Applies': '{:.0f}', 'Mean Applies': '{:.1f}',
                'Apply Rate %': '{:.1f}%', 'Avg Days Live': '{:.0f}'
            })
            styled = apply_manual_gradient(styled, 'Median Clicks', 'green')
            styled = apply_manual_gradient(styled, 'Median Applies', 'blue')
            styled = apply_manual_gradient(styled, 'Apply Rate %', 'orange')

            st.dataframe(styled, width='stretch', hide_index=True,
                         height=min(600, 35 * len(sector) + 40))
        else:
            st.info("Not enough data (need 5+ vacancies per occupation).")
    else:
        st.info("No data available for benchmarks.")


# ============================================================================
# TAB 6: LAUNCH TIMING
# ============================================================================

def create_launch_timing_tab(launch_df):
    """Analyse vacancy performance by day offset from launch, with day-of-week analysis."""

    st.subheader("Vacancy Performance by Day Since Launch")
    st.caption("Each vacancy is normalised to Day 0 (first event). Averages across all vacancies show the typical performance curve.")

    if launch_df is None or len(launch_df) == 0:
        st.warning("No launch timing data available. This requires event-level data from the enriched table.")
        return

    df = launch_df.copy()

    # Add occupation column for filtering
    if 'occupational_fields' in df.columns:
        df['occupation'] = df['occupational_fields'].apply(lambda x:
            str(x).split('|')[0].strip() if pd.notna(x) and str(x).strip() else 'Unknown'
        )

    # Filters
    col1, col2 = st.columns(2)
    with col1:
        if 'occupation' in df.columns:
            occupations = sorted(df['occupation'].dropna().unique())
            selected_occs = st.multiselect("Filter by Occupation", occupations, key='launch_occ')
            if selected_occs:
                df = df[df['occupation'].isin(selected_occs)]
    with col2:
        metric_choice = st.radio("Metric", ["Clicks (Views)", "Applies"], horizontal=True, key='launch_metric')
        metric_col = 'clicks' if 'Clicks' in metric_choice else 'applies'

    vacancy_count = df['entity_id_str'].nunique()
    st.info(f"Analysing **{vacancy_count:,}** vacancies")

    # Map BigQuery DAYOFWEEK (1=Sun, 2=Mon, ... 7=Sat) to readable names
    dow_map = {1: 'Sunday', 2: 'Monday', 3: 'Tuesday', 4: 'Wednesday',
               5: 'Thursday', 6: 'Friday', 7: 'Saturday'}
    dow_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
    df['launch_day_name'] = df['launch_dow'].map(dow_map)
    df['event_day_name'] = df['event_dow'].map(dow_map)

    # ------------------------------------------------------------------
    # Chart 1: Average performance by day offset (Day 0 - Day 30)
    # ------------------------------------------------------------------
    st.markdown("### Performance Curve (Day 0 — Day 30)")

    avg_by_day = df.groupby('day_offset')[metric_col].mean().reset_index()
    avg_by_day.columns = ['Day', metric_col]

    fig1 = go.Figure()
    fig1.add_trace(go.Scatter(
        x=avg_by_day['Day'], y=avg_by_day[metric_col],
        mode='lines+markers',
        name=f'Avg {metric_col.title()}',
        line=dict(color='#2563eb', width=2),
        marker=dict(size=5)
    ))

    # Shade weekends (approximate: if Day 0 is any day, weekends fall at offsets
    # that depend on launch day — so instead shade typical weekend bands)
    # We'll use the event_dow to mark which day_offsets tend to be weekends
    weekend_offsets = df[df['event_day_name'].isin(['Saturday', 'Sunday'])].groupby('day_offset').size()
    all_offsets = df.groupby('day_offset').size()
    weekend_pct = (weekend_offsets / all_offsets).fillna(0)
    # Shade offsets where >40% of events fall on weekends
    weekend_days = weekend_pct[weekend_pct > 0.4].index.tolist()

    for d in weekend_days:
        fig1.add_vrect(x0=d - 0.5, x1=d + 0.5, fillcolor='rgba(200,200,200,0.2)',
                       line_width=0, layer='below')

    fig1.update_layout(
        height=400,
        xaxis_title='Days Since Launch',
        yaxis_title=f'Avg {metric_col.title()} per Vacancy',
        hovermode='x unified',
        annotations=[dict(x=0.98, y=0.98, xref='paper', yref='paper',
                          text='Shaded = Weekend', showarrow=False,
                          font=dict(size=10, color='grey'))]
    )
    st.plotly_chart(fig1, width='stretch')

    # ------------------------------------------------------------------
    # Chart 2: Performance curve by launch day of week
    # ------------------------------------------------------------------
    st.markdown("### Performance by Launch Day of Week")
    st.caption("Which day of the week is best to launch a vacancy?")

    avg_by_dow = df.groupby(['launch_day_name', 'day_offset'])[metric_col].mean().reset_index()

    colors = {'Monday': '#2563eb', 'Tuesday': '#7c3aed', 'Wednesday': '#059669',
              'Thursday': '#d97706', 'Friday': '#dc2626', 'Saturday': '#6b7280', 'Sunday': '#9ca3af'}

    fig2 = go.Figure()
    for day in dow_order:
        day_data = avg_by_dow[avg_by_dow['launch_day_name'] == day]
        if len(day_data) > 0:
            fig2.add_trace(go.Scatter(
                x=day_data['day_offset'], y=day_data[metric_col],
                mode='lines',
                name=day,
                line=dict(color=colors.get(day, '#333'), width=2),
                opacity=0.8
            ))

    fig2.update_layout(
        height=450,
        xaxis_title='Days Since Launch',
        yaxis_title=f'Avg {metric_col.title()} per Vacancy',
        hovermode='x unified',
        legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='right', x=1)
    )
    st.plotly_chart(fig2, width='stretch')

    # ------------------------------------------------------------------
    # Summary table: first 7 days total by launch day
    # ------------------------------------------------------------------
    st.markdown("### First 7 Days Summary by Launch Day")

    first_week = df[df['day_offset'] <= 6].copy()
    summary = first_week.groupby('launch_day_name').agg(
        vacancies=('entity_id_str', 'nunique'),
        avg_clicks_7d=('clicks', 'mean'),
        avg_applies_7d=('applies', 'mean')
    ).reset_index()

    # Multiply by 7 to get total over first week (avg per day × 7 days)
    summary['avg_clicks_7d'] = (summary['avg_clicks_7d'] * 7).round(1)
    summary['avg_applies_7d'] = (summary['avg_applies_7d'] * 7).round(1)

    # Sort by day of week
    summary['sort_key'] = summary['launch_day_name'].map({d: i for i, d in enumerate(dow_order)})
    summary = summary.sort_values('sort_key').drop(columns='sort_key')

    summary.columns = ['Launch Day', 'Vacancies', 'Avg Clicks (First 7 Days)', 'Avg Applies (First 7 Days)']

    def _manual_grad(df_styled, col, r, g, b):
        col_data = summary[col].dropna()
        if len(col_data) == 0:
            return df_styled
        cmin, cmax = col_data.min(), col_data.max()

        def _safe_grad(v):
            try:
                if pd.isna(v) or cmax <= cmin:
                    return ''
                intensity = min((float(v) - float(cmin)) / (float(cmax) - float(cmin)), 1) * 0.3
                return f'background-color: rgba({r},{g},{b},{intensity:.2f})'
            except (TypeError, ValueError):
                return ''

        return df_styled.map(_safe_grad, subset=[col])

    styled = summary.style.format({'Avg Clicks (First 7 Days)': '{:.1f}', 'Avg Applies (First 7 Days)': '{:.1f}'})
    styled = _manual_grad(styled, 'Avg Clicks (First 7 Days)', 34, 139, 34)
    styled = _manual_grad(styled, 'Avg Applies (First 7 Days)', 30, 90, 200)
    st.dataframe(styled, width='stretch', hide_index=True)


def generate_section_commentary(section, data):
    """Generate template-based commentary for each report section.

    Args:
        section: One of 'scatter', 'benchmark', 'postings', 'roi', 'media'.
        data: Dict with relevant metrics for the section.

    Returns:
        Markdown-formatted string with 2-4 sentences of data-driven insight.
    """
    fallback = "Insufficient data for detailed commentary."

    if section == 'scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        worst_performers = data.get('worst_performers', [])

        if total == 0:
            return fallback
        if total == 1:
            return ("Only one vacancy is available for analysis, which limits "
                    "benchmarking comparisons. A larger sample of roles will "
                    "enable more meaningful performance insights.")

        parts = []
        parts.append(
            f"Of **{total}** vacancies, **{benchmarkable}** "
            f"({'all' if benchmarkable == total else f'{benchmarkable / total:.0%}'}) "
            f"have sufficient market data for benchmarking"
            f"{f', while **{no_benchmark}** lack benchmark data due to low sample sizes in their occupation category' if no_benchmark > 0 else ''}."
        )

        if top_performers:
            top = top_performers[0]
            parts.append(
                f"**{top['title']}** ({top['occupation']}) is a standout performer, "
                f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views "
                f"and {top['applies_diff_pct']:+.0f}% on applies."
            )

        if zero_applies > 0:
            pct_zero = zero_applies / total * 100
            parts.append(
                f"**{zero_applies}** role{'s' if zero_applies != 1 else ''} "
                f"({pct_zero:.0f}% of total) received zero apply clicks and "
                f"may benefit from revised job descriptions or enhanced visibility."
            )

        return " ".join(parts)

    elif section == 'benchmark':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return fallback

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        views_word = "more" if views_diff >= 0 else "fewer"

        parts = []
        parts.append(
            f"Your vacancies received **{abs(views_diff):.0f}% {views_word} views** "
            f"than the market average, "
            f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}."
        )

        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            parts.append(
                f"Apply rates sit **{abs(applies_diff):.0f}% {applies_word} benchmark**"
                f"{', reflecting strong candidate engagement' if applies_diff >= 10 else ', suggesting job descriptions may benefit from enhancement' if applies_diff < -10 else ', broadly in line with market expectations'}."
            )

        parts.append(
            f"This analysis is based on **{num_jobs}** vacancies posted by {client_name} "
            f"during the report period."
        )
        return " ".join(parts)

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0:
            return fallback
        if by_type is None or len(by_type) == 0:
            return f"{client_name} posted **{total_jobs}** vacancies with **{total_applies:,}** total apply clicks during this period."

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        parts = []

        if len(by_type) == 1:
            parts.append(
                f"All **{total_jobs}** vacancies fall under **{top['occupation']}**, "
                f"generating **{int(top['apply_clicks']):,}** apply clicks."
            )
        else:
            parts.append(
                f"**{top['occupation']}** leads apply generation with "
                f"**{int(top['apply_clicks']):,}** apply clicks across "
                f"**{int(top['jobs_posted'])}** postings."
            )
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                parts.append(
                    f"**{bottom['occupation']}** received no apply clicks despite "
                    f"**{int(bottom['jobs_posted'])}** postings — these roles may "
                    f"benefit from revised job titles or enhanced descriptions."
                )
            elif len(by_type) > 1:
                parts.append(
                    f"**{bottom['occupation']}** generated the fewest applies "
                    f"(**{int(bottom['apply_clicks']):,}**) and may warrant "
                    f"targeted improvements to boost candidate engagement."
                )

        if total_applies == 0:
            parts.append(
                "No apply clicks were recorded across any category — reviewing "
                "listing quality and distribution channels is recommended."
            )
        return " ".join(parts)

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return "Enter your annual spend and rate card price to generate ROI commentary."

        parts = []
        if saving_pct > 0:
            parts.append(
                f"Your advertising investment delivers a **{saving_pct:.0f}% saving** "
                f"compared to rate card pricing, demonstrating strong value from "
                f"the platform partnership."
            )
        else:
            parts.append(
                "Current spend exceeds rate card value — reviewing the pricing "
                "structure or consolidating lower-performing listings may improve "
                "overall return on investment."
            )

        total_applies = data.get('total_applies', 0)
        if total_applies > 0:
            parts.append(
                f"At **\u00a3{cost_per_apply:,.2f} per apply**, each candidate "
                f"enquiry represents a cost-effective acquisition channel."
            )

        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            parts.append(
                f"**{best['occupation']}** achieves the best cost efficiency at "
                f"\u00a3{best['cost_per_apply']:,.2f} per apply, while "
                f"**{worst['occupation']}** is the most expensive at "
                f"\u00a3{worst['cost_per_apply']:,.2f}."
            )

        return " ".join(parts)

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return "Media source data is not yet available for this client."

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        parts = []

        parts.append(
            f"**{top_source['source_category']}** is the leading traffic source, "
            f"generating **{int(top_source['total_applies']):,}** apply clicks "
            f"from **{int(top_source['total_clicks']):,}** views."
        )

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        if best_conv['conversion_rate'] > 0:
            parts.append(
                f"**{best_conv['source_category']}** achieves the highest "
                f"view-to-apply conversion rate at **{best_conv['conversion_rate']:.1f}%**."
            )

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            parts.append(
                f"Paid channels ({paid['source_category']}) contributed "
                f"**{int(paid['total_applies']):,}** applies with a "
                f"**{paid['conversion_rate']:.1f}%** conversion rate."
            )

        return " ".join(parts)

    return fallback


def generate_section_commentary_structured(section, data):
    """Structured commentary for PPTX template — returns dict with intro + bullet points.

    Returns: dict with keys 'intro', 'point_1', 'point_2', 'point_3' (last is optional).
    Each value is plain text (no markdown). Empty string for unused points.
    """
    def _clean(text):
        """Strip markdown bold markers."""
        import re
        return re.sub(r'\*\*(.+?)\*\*', r'\1', text or '')

    if section == 'benchmark_scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        client_name = data.get('client_name', 'This client')

        if total == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': '', 'point_3': ''}

        intro = (f"{client_name} posted {total} vacancies during this period. "
                 f"Of these, {benchmarkable} have sufficient market data to benchmark against comparable public sector roles.")

        point_1 = ''
        if top_performers:
            top = top_performers[0]
            point_1 = (f"{top['title']} ({top['occupation']}) is the standout performer, "
                       f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views and "
                       f"{top['applies_diff_pct']:+.0f}% on applies.")

        point_2 = ''
        if zero_applies > 0:
            pct = zero_applies / total * 100 if total > 0 else 0
            point_2 = (f"{zero_applies} role{'s' if zero_applies != 1 else ''} ({pct:.0f}% of total) "
                       f"received zero apply clicks — these may benefit from revised job descriptions or enhanced visibility.")

        point_3 = ''
        if no_benchmark > 0:
            point_3 = (f"{no_benchmark} role{'s' if no_benchmark != 1 else ''} could not be benchmarked "
                       f"due to low market sample sizes in their occupation category.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    elif section == 'benchmark_average':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        intro = (f"Across {num_jobs} vacancies, {client_name} averaged {client_clicks:,.0f} views "
                 f"and {client_applies:,.1f} applies per role.")

        views_word = "more" if views_diff >= 0 else "fewer"
        point_1 = (f"Your vacancies received {abs(views_diff):.0f}% {views_word} views than the market average — "
                   f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}.")

        point_2 = ''
        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            tone = ('reflecting strong candidate engagement' if applies_diff >= 10
                    else 'suggesting job descriptions may benefit from enhancement' if applies_diff < -10
                    else 'broadly in line with market expectations')
            point_2 = f"Apply rates sit {abs(applies_diff):.0f}% {applies_word} benchmark, {tone}."

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0 or by_type is None or len(by_type) == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        intro = (f"{client_name} posted {total_jobs} vacancies generating {total_applies:,} apply clicks across "
                 f"{len(by_type)} occupation categories.")

        point_1 = (f"{top['occupation']} leads apply generation with {int(top['apply_clicks']):,} apply clicks across "
                   f"{int(top['jobs_posted'])} postings — your strongest performing category.")

        point_2 = ''
        if len(by_type) > 1:
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                point_2 = (f"{bottom['occupation']} received no apply clicks despite {int(bottom['jobs_posted'])} postings — "
                           f"these roles may benefit from revised titles or enhanced descriptions.")
            else:
                point_2 = (f"{bottom['occupation']} generated the fewest applies ({int(bottom['apply_clicks']):,}) "
                           f"and may warrant targeted improvements.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        rate_card_total = data.get('rate_card_total', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return {'intro': 'Enter your annual spend and rate card price to generate ROI commentary.',
                    'point_1': '', 'point_2': ''}

        saving_amount = rate_card_total - annual_spend
        intro = (f"Across {num_jobs} vacancies, your subscription delivered £{saving_amount:,.0f} of value "
                 f"compared to rate card pricing — a {saving_pct:.0f}% saving.")

        point_1 = (f"At £{cost_per_apply:,.2f} per apply, each candidate enquiry represents a "
                   f"cost-effective acquisition channel for {data.get('client_name', 'your team')}.")

        point_2 = ''
        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            point_2 = (f"{best['occupation']} achieves the best cost efficiency at £{best['cost_per_apply']:,.2f} per apply, "
                       f"while {worst['occupation']} is the most expensive at £{worst['cost_per_apply']:,.2f}.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return {'intro': 'Media source data is not yet available for this client.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        intro = (f"{client_name}'s vacancies received traffic from {len(cat_stats)} distinct channels. "
                 f"{top_source['source_category']} is the leading source, generating "
                 f"{int(top_source['total_applies']):,} applies from {int(top_source['total_clicks']):,} views.")

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        point_1 = ''
        if best_conv['conversion_rate'] > 0:
            point_1 = (f"{best_conv['source_category']} achieves the highest view-to-apply conversion rate at "
                       f"{best_conv['conversion_rate']:.1f}%, indicating well-matched candidates from this channel.")

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        point_2 = ''
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            point_2 = (f"Paid channels ({paid['source_category']}) contributed {int(paid['total_applies']):,} applies "
                       f"with a {paid['conversion_rate']:.1f}% conversion rate.")

        point_3 = ''
        if len(cat_stats) > 1:
            second_source = sorted_stats.iloc[1] if len(sorted_stats) > 1 else None
            if second_source is not None:
                point_3 = (f"{second_source['source_category']} is the second strongest channel with "
                           f"{int(second_source['total_applies']):,} applies — providing diversified candidate flow.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    return {'intro': 'No commentary available.', 'point_1': '', 'point_2': '', 'point_3': ''}


# ============================================================================
# CLIENT REPORT TAB
# ============================================================================

def create_client_report_tab(df, media_df=None):
    """Create the Client Report tab — replicates the branded PDF advertising report."""
    st.header("Client Advertising Report")

    # --- Controls ---
    col_client, col_dates = st.columns(2)
    with col_client:
        # Client selector — uses organization_name (the actual client/employer)
        org_col = 'organization_name' if 'organization_name' in df.columns else 'importer_name'
        orgs = sorted(df[org_col].dropna().unique())
        # Filter out blanks and unknowns
        orgs = [o for o in orgs if o and str(o).strip() not in ('', 'Unknown', 'nan')]
        org_counts = df.groupby(org_col).size().to_dict()
        org_labels = [f"{name} ({org_counts.get(name, 0):,} vacancies)" for name in orgs]
        selected_idx = st.selectbox(
            "Select Client / Organisation", range(len(orgs)),
            format_func=lambda i: org_labels[i],
            key='report_client'
        )
        selected_client = orgs[selected_idx]
    with col_dates:
        min_date = df['first_event_date'].dropna().min()
        max_date = df['last_event_date'].dropna().max()
        if pd.notna(min_date) and pd.notna(max_date):
            min_d = min_date.date() if hasattr(min_date, 'date') else min_date
            max_d = max_date.date() if hasattr(max_date, 'date') else max_date
        else:
            min_d = datetime.now().date() - timedelta(days=365)
            max_d = datetime.now().date()
        report_dates = st.date_input("Report Period", [min_d, max_d], key='report_dates')

    with st.expander("Cost & Report Settings", expanded=False):
        cost_col1, cost_col2 = st.columns(2)
        with cost_col1:
            annual_spend = st.number_input(
                "Annual Spend (GBP)", value=0.0, step=100.0, format="%.2f",
                key='report_spend',
                help="Enter 0 to skip the ROI section"
            )
        with cost_col2:
            rate_card_price = st.number_input(
                "Rate Card Price per Job (GBP)", value=600.0, step=10.0, format="%.2f",
                key='report_rate_card'
            )
        settings_col1, settings_col2 = st.columns(2)
        with settings_col1:
            include_self = st.checkbox(
                "Include self in benchmark",
                value=False,
                key='report_include_self',
                help="When unchecked, the selected client is excluded from the benchmark average (recommended for fair comparison)"
            )
        with settings_col2:
            st.markdown("**Contact details** (for PDF)")
            contact_name = st.text_input("Account Manager", key='report_contact_name', placeholder="e.g. Jane Smith")
            contact_title = st.text_input("Title", key='report_contact_title', placeholder="e.g. Account Director")
            contact_email = st.text_input("Email", key='report_contact_email', placeholder="e.g. jane@jgp.co.uk")
            contact_phone = st.text_input("Phone", key='report_contact_phone', placeholder="e.g. 020 7946 0958")

    generate_clicked = st.button("Generate Report", type="primary", key='report_generate')

    if not generate_clicked and 'report_generated' not in st.session_state:
        st.info("Select a client and click **Generate Report** to build the advertising report.")
        return

    st.session_state['report_generated'] = True

    # --- Data preparation ---
    if len(report_dates) < 2:
        st.warning("Please select a start and end date.")
        return

    report_start, report_end = report_dates[0], report_dates[1]

    # Client data — filter on organization_name (or importer_name fallback)
    client_df = df[df[org_col] == selected_client].copy()
    client_df = client_df[
        (client_df['last_event_date'].dt.date >= report_start) &
        (client_df['first_event_date'].dt.date <= report_end)
    ]

    if len(client_df) == 0:
        st.warning(f"No vacancies found for **{selected_client}** in the selected date range.")
        return

    # Benchmark = market data in same date range (exclude self by default for fair comparison)
    bench_mask = (
        (df['last_event_date'].dt.date >= report_start) &
        (df['first_event_date'].dt.date <= report_end)
    )
    if not include_self:
        bench_mask = bench_mask & (df[org_col] != selected_client)
    benchmark_df = df[bench_mask].copy()

    # Media data for client — match via entity_id (most reliable link)
    # Note: prepare_enriched_data() renames entity_id_str → entity_id in df,
    # but media_df keeps the original entity_id_str column name.
    client_media = None
    if media_df is not None and len(media_df) > 0:
        # Determine the entity ID column name in each dataframe
        client_eid_col = 'entity_id' if 'entity_id' in client_df.columns else 'entity_id_str'
        media_eid_col = 'entity_id_str' if 'entity_id_str' in media_df.columns else 'entity_id'
        if client_eid_col in client_df.columns and media_eid_col in media_df.columns:
            client_entity_ids = client_df[client_eid_col].dropna().unique()
            client_media = media_df[media_df[media_eid_col].isin(client_entity_ids)].copy()
        if client_media is None or len(client_media) == 0:
            client_media = None
        else:
            client_media = apply_media_categories(client_media)

    # Store all figures for PDF export
    report_figures = {}

    st.markdown("---")

    # ===================================================================
    # SECTION 1: BENCHMARKING SCATTER
    # ===================================================================
    st.subheader("Benchmarking Jobs")

    # Calculate per-occupation benchmark averages (from ALL clients)
    occ_benchmarks = benchmark_df.groupby('occupation').agg(
        avg_clicks=('clicks', 'mean'),
        avg_applies=('applies', 'mean'),
        vacancy_count=('clicks', 'count')
    ).reset_index()

    # Only use occupations with enough data for reliable benchmarks
    MIN_BENCHMARK_VACANCIES = 5
    reliable_occs = occ_benchmarks[occ_benchmarks['vacancy_count'] >= MIN_BENCHMARK_VACANCIES]

    # Calculate % difference from benchmark for each client vacancy (vectorized)
    scatter_df = client_df[['title', 'occupation', 'clicks', 'applies']].copy()
    scatter_df = scatter_df.merge(
        reliable_occs[['occupation', 'avg_clicks', 'avg_applies']],
        on='occupation', how='left'
    )

    # Vectorized % diff calculations (safe division)
    scatter_df['views_diff_pct'] = (
        (scatter_df['clicks'] - scatter_df['avg_clicks'])
        / scatter_df['avg_clicks'].replace(0, np.nan) * 100
    ).fillna(0)
    scatter_df['applies_diff_pct'] = (
        (scatter_df['applies'] - scatter_df['avg_applies'])
        / scatter_df['avg_applies'].replace(0, np.nan) * 100
    ).fillna(0)

    # Categorize vacancies into 4 groups (matching original PDF categories)
    has_benchmark = scatter_df['avg_clicks'].notna()
    zero_applies = scatter_df['applies'] == 0
    occ_median_clicks = scatter_df.loc[has_benchmark, 'clicks'].median() if has_benchmark.any() else 0
    high_traffic = scatter_df['clicks'] > occ_median_clicks

    scatter_df['category'] = np.select(
        [has_benchmark & ~zero_applies,
         has_benchmark & zero_applies & high_traffic,
         has_benchmark & zero_applies & ~high_traffic],
        ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)'],
        default='Low Sample (No Benchmark)'
    )

    # Clean up merge columns and fill NaN diffs for no-benchmark rows
    scatter_df.drop(columns=['avg_clicks', 'avg_applies'], inplace=True)
    scatter_df.loc[~has_benchmark, ['views_diff_pct', 'applies_diff_pct']] = 0

    # --- Brand colour + marker maps for 4 categories ---
    from theme.colors import JGP_COLORS
    category_colors = {
        'Benchmarkable': JGP_COLORS['primary'],
        'Zero Applies (Possible Redirect)': JGP_COLORS['amber'],
        'Zero Applies (Low Traffic)': JGP_COLORS['negative'],
        'Low Sample (No Benchmark)': JGP_COLORS['light_purple'],
    }
    category_symbols = {
        'Benchmarkable': 'triangle-up',
        'Zero Applies (Possible Redirect)': 'diamond',
        'Zero Applies (Low Traffic)': 'x',
        'Low Sample (No Benchmark)': 'circle',
    }

    chart_col, commentary_col = st.columns([3, 2])

    with chart_col:
        if len(scatter_df) > 0:
            fig_scatter = go.Figure()
            for cat in ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)', 'Low Sample (No Benchmark)']:
                cat_data = scatter_df[scatter_df['category'] == cat]
                if len(cat_data) == 0:
                    continue
                fig_scatter.add_trace(go.Scatter(
                    x=cat_data['applies_diff_pct'],
                    y=cat_data['views_diff_pct'],
                    mode='markers',
                    name=cat,
                    marker=dict(
                        color=category_colors[cat],
                        symbol=category_symbols[cat],
                        size=10,
                        line=dict(width=1, color='white')
                    ),
                    customdata=cat_data[['title', 'occupation', 'clicks', 'applies']].values,
                    hovertemplate=(
                        "<b>%{customdata[0]}</b><br>"
                        "Occupation: %{customdata[1]}<br>"
                        "Views: %{customdata[2]}<br>"
                        "Applies: %{customdata[3]}<br>"
                        "Views vs Bench: %{y:+.0f}%<br>"
                        "Applies vs Bench: %{x:+.0f}%<extra></extra>"
                    ),
                ))
            fig_scatter.add_hline(y=0, line_dash="dash", line_color="grey", opacity=0.5)
            fig_scatter.add_vline(x=0, line_dash="dash", line_color="grey", opacity=0.5)
            fig_scatter.update_layout(
                height=500,
                showlegend=True,
                legend=dict(orientation='h', y=-0.15, font=dict(size=11)),
                xaxis_title="Applies Difference from Benchmark (%)",
                yaxis_title="Views Difference from Benchmark (%)",
                plot_bgcolor='rgba(0,0,0,0)',
                xaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
                yaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
            )
            st.plotly_chart(fig_scatter, use_container_width=True)
            report_figures['scatter'] = fig_scatter

    with commentary_col:
        benchmarkable_count = len(scatter_df[scatter_df['category'] == 'Benchmarkable'])
        zero_redirect = len(scatter_df[scatter_df['category'] == 'Zero Applies (Possible Redirect)'])
        zero_low = len(scatter_df[scatter_df['category'] == 'Zero Applies (Low Traffic)'])
        no_bench_count = len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)'])

        kpi_c1, kpi_c2 = st.columns(2)
        with kpi_c1:
            st.metric("Benchmarkable", benchmarkable_count)
            st.metric("Zero Applies (Redirect)", zero_redirect)
        with kpi_c2:
            st.metric("Zero Applies (Low Traffic)", zero_low)
            st.metric("No Benchmark", no_bench_count)

        # Build top/worst performers for commentary
        benchmarkable_rows = scatter_df[scatter_df['category'] == 'Benchmarkable'].copy()
        top_performers = []
        worst_performers = []
        if len(benchmarkable_rows) > 0:
            benchmarkable_rows['_score'] = benchmarkable_rows['views_diff_pct'] + benchmarkable_rows['applies_diff_pct']
            top_sorted = benchmarkable_rows.nlargest(3, '_score')
            top_performers = top_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')
            worst_sorted = benchmarkable_rows.nsmallest(3, '_score')
            worst_performers = worst_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')

        commentary = generate_section_commentary('scatter', {
            'total_count': len(scatter_df),
            'benchmarkable_count': benchmarkable_count,
            'zero_applies_count': zero_redirect + zero_low,
            'no_benchmark_count': no_bench_count,
            'top_performers': top_performers,
            'worst_performers': worst_performers,
        })
        st.markdown(commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 2: BENCHMARKING SUMMARY
    # ===================================================================
    st.subheader("Benchmarking Summary")

    benchmark_avg_clicks = benchmark_df['clicks'].mean() if len(benchmark_df) > 0 else 0
    benchmark_avg_applies = benchmark_df['applies'].mean() if len(benchmark_df) > 0 else 0
    client_avg_clicks = client_df['clicks'].mean() if len(client_df) > 0 else 0
    client_avg_applies = client_df['applies'].mean() if len(client_df) > 0 else 0

    kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
    with kpi_col1:
        st.metric("Benchmark Avg. Views", f"{benchmark_avg_clicks:,.0f}")
    with kpi_col2:
        views_delta = client_avg_clicks - benchmark_avg_clicks
        st.metric("Your Jobs - Avg. Views", f"{client_avg_clicks:,.0f}",
                  delta=f"{views_delta:+,.0f} vs benchmark")
    with kpi_col3:
        st.metric("Benchmark Avg. Applies", f"{benchmark_avg_applies:,.1f}")
    with kpi_col4:
        applies_delta = client_avg_applies - benchmark_avg_applies
        st.metric("Your Jobs - Avg. Applies", f"{client_avg_applies:,.1f}",
                  delta=f"{applies_delta:+,.1f} vs benchmark")

    # Bar charts (60%) + commentary (40%)
    bench_chart_col, bench_commentary_col = st.columns([3, 2])

    with bench_chart_col:
        views_pct = (client_avg_clicks / benchmark_avg_clicks * 100) if benchmark_avg_clicks > 0 else 0
        applies_pct = (client_avg_applies / benchmark_avg_applies * 100) if benchmark_avg_applies > 0 else 0

        fig_bench = go.Figure()
        fig_bench.add_trace(go.Bar(
            x=['Views', 'Applies'],
            y=[views_pct, applies_pct],
            marker_color=[
                JGP_COLORS['positive'] if views_pct >= 100 else JGP_COLORS['amber'],
                JGP_COLORS['positive'] if applies_pct >= 100 else JGP_COLORS['amber'],
            ],
            text=[f"{views_pct:.0f}%", f"{applies_pct:.0f}%"],
            textposition='outside',
            textfont=dict(size=14, color=JGP_COLORS['deep_blue']),
        ))
        fig_bench.add_hline(
            y=100, line_dash="dash", line_width=3, line_color=JGP_COLORS['accent'],
            annotation_text="Benchmark (100%)",
            annotation_position="top right",
            annotation_bgcolor="white",
            annotation_bordercolor=JGP_COLORS['deep_blue'],
            annotation_borderwidth=1,
            annotation_font=dict(color=JGP_COLORS['deep_blue'], size=14),
        )
        fig_bench.update_layout(
            title="Your Performance vs Market Benchmark",
            yaxis_title="% of Benchmark",
            height=400, showlegend=False,
            yaxis_range=[0, max(views_pct, applies_pct, 100) * 1.25],
            plot_bgcolor='rgba(0,0,0,0)',
        )
        st.plotly_chart(fig_bench, use_container_width=True)
        report_figures['benchmark_combined'] = fig_bench

    with bench_commentary_col:
        bench_commentary = generate_section_commentary('benchmark', {
            'client_avg_clicks': client_avg_clicks,
            'benchmark_avg_clicks': benchmark_avg_clicks,
            'client_avg_applies': client_avg_applies,
            'benchmark_avg_applies': benchmark_avg_applies,
            'num_jobs': len(client_df),
            'client_name': selected_client,
        })
        st.markdown(bench_commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 3: JOB POSTINGS BY TYPE
    # ===================================================================
    st.subheader("Job Postings by Type")

    by_type = client_df.groupby('occupation').agg(
        jobs_posted=('clicks', 'count'),
        apply_clicks=('applies', 'sum')
    ).reset_index().sort_values('jobs_posted', ascending=True)

    by_type = by_type[by_type['jobs_posted'] >= 1]

    chart_col3, commentary_col3 = st.columns([3, 2])

    with chart_col3:
        fig_postings = go.Figure()
        fig_postings.add_trace(go.Bar(
            y=by_type['occupation'], x=by_type['jobs_posted'],
            name='Jobs Posted', orientation='h',
            marker_color=JGP_COLORS['primary'],
            text=by_type['jobs_posted'], textposition='outside'
        ))
        fig_postings.add_trace(go.Bar(
            y=by_type['occupation'], x=by_type['apply_clicks'],
            name='Apply Clicks', orientation='h',
            marker_color=JGP_COLORS['accent'],
            text=by_type['apply_clicks'].astype(int), textposition='outside',
            textfont=dict(color=JGP_COLORS['deep_blue']),
        ))
        fig_postings.update_layout(
            barmode='group', height=max(400, len(by_type) * 40),
            legend=dict(orientation='h', y=-0.1),
            xaxis_title="Count", yaxis_title="",
            plot_bgcolor='rgba(0,0,0,0)',
            bargap=0.1, bargroupgap=0.0,
        )
        st.plotly_chart(fig_postings, use_container_width=True)
        report_figures['postings'] = fig_postings

    with commentary_col3:
        total_jobs = len(client_df)
        total_applies_val = int(client_df['applies'].sum())
        postings_commentary = generate_section_commentary('postings', {
            'total_jobs': total_jobs,
            'total_applies': total_applies_val,
            'by_type': by_type,
            'client_name': selected_client,
        })
        st.markdown(postings_commentary)

    st.markdown("---")

    # ===================================================================
    # SECTION 4: ADVERTISING ROI
    # ===================================================================
    st.subheader("Advertising ROI")

    num_jobs = len(client_df)
    total_clicks = int(client_df['clicks'].sum())
    total_applies_val = int(client_df['applies'].sum())

    if annual_spend > 0:
        cost_per_job = annual_spend / num_jobs if num_jobs > 0 else 0
        cost_per_view = annual_spend / total_clicks if total_clicks > 0 else 0
        cost_per_apply = annual_spend / total_applies_val if total_applies_val > 0 else 0

        roi_kpi1, roi_kpi2, roi_kpi3, roi_kpi4 = st.columns(4)
        with roi_kpi1:
            st.metric("Jobs Advertised", f"{num_jobs:,}")
        with roi_kpi2:
            st.metric("Cost per Job", f"£{cost_per_job:,.2f}")
        with roi_kpi3:
            st.metric("Cost per View", f"£{cost_per_view:,.2f}")
        with roi_kpi4:
            st.metric("Cost per Apply", f"£{cost_per_apply:,.2f}")

        roi_chart_col, roi_commentary_col = st.columns([3, 2])

        with roi_chart_col:
            rate_card_total = rate_card_price * num_jobs
            saving_pct = ((rate_card_total - annual_spend) / rate_card_total * 100) if rate_card_total > 0 else 0

            fig_roi = go.Figure()
            fig_roi.add_trace(go.Bar(
                x=['Your Spend'], y=[annual_spend],
                name='Your Spend', marker_color=JGP_COLORS['primary'],
                text=[f"£{annual_spend:,.0f}"], textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_roi.add_trace(go.Bar(
                x=['Rate Card Value'], y=[rate_card_total],
                name='Rate Card Value', marker_color=JGP_COLORS['deep_blue'],
                text=[f"£{rate_card_total:,.0f}"], textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_roi.update_layout(
                title=f"Cost vs Rate Card (Saving: {saving_pct:.0f}%)",
                height=350, showlegend=True, yaxis_title="GBP",
                plot_bgcolor='rgba(0,0,0,0)',
            )
            st.plotly_chart(fig_roi, use_container_width=True)
            report_figures['roi_cost'] = fig_roi

            # Cost per apply by type chart
            roi_by_type = client_df.groupby('occupation').agg(
                total_applies=('applies', 'sum'),
                job_count=('clicks', 'count')
            ).reset_index()
            roi_by_type = roi_by_type[roi_by_type['total_applies'] > 0]
            roi_by_type['cost_allocated'] = annual_spend * (roi_by_type['job_count'] / roi_by_type['job_count'].sum())
            roi_by_type['cost_per_apply'] = roi_by_type['cost_allocated'] / roi_by_type['total_applies']
            roi_by_type = roi_by_type.sort_values('cost_per_apply', ascending=True)

            if len(roi_by_type) > 0:
                fig_cpa = go.Figure()
                fig_cpa.add_trace(go.Bar(
                    y=roi_by_type['occupation'], x=roi_by_type['cost_per_apply'],
                    orientation='h', marker_color=JGP_COLORS['amber'],
                    text=roi_by_type['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                    textposition='outside'
                ))
                fig_cpa.update_layout(
                    title="Cost per Apply by Job Type",
                    height=max(300, len(roi_by_type) * 35),
                    xaxis_title="Cost per Apply (GBP)", yaxis_title="",
                    plot_bgcolor='rgba(0,0,0,0)',
                )
                st.plotly_chart(fig_cpa, use_container_width=True)
                report_figures['roi_cpa'] = fig_cpa

        with roi_commentary_col:
            roi_commentary = generate_section_commentary('roi', {
                'annual_spend': annual_spend,
                'rate_card_price': rate_card_price,
                'num_jobs': num_jobs,
                'total_clicks': total_clicks,
                'total_applies': total_applies_val,
                'cost_per_job': cost_per_job,
                'cost_per_view': cost_per_view,
                'cost_per_apply': cost_per_apply,
                'saving_pct': saving_pct,
                'roi_by_type': roi_by_type if len(roi_by_type) > 0 else None,
            })
            st.markdown(roi_commentary)
    else:
        st.metric("Jobs Advertised", f"{num_jobs:,}")
        st.info("Enter your **Annual Spend** and **Rate Card Price** in the Cost & Report Settings above to see ROI analysis.")

    st.markdown("---")

    # ===================================================================
    # SECTION 5: MEDIA PERFORMANCE
    # ===================================================================
    st.subheader(f"Media Performance — {selected_client}")

    cat_stats = None  # Initialise before conditional block so it's in scope for PDF commentary
    if client_media is not None and len(client_media) > 0:
        media_vac_count = client_media['entity_id_str'].nunique() if 'entity_id_str' in client_media.columns else len(client_media)
        st.caption(f"Showing media data for **{media_vac_count:,}** vacancies belonging to {selected_client}")
        # Category-level summary
        cat_stats = client_media.groupby('source_category').agg(
            total_clicks=('clicks', 'sum'),
            total_applies=('applies', 'sum'),
            vacancy_count=('entity_id_str', 'nunique')
        ).reset_index()
        cat_stats['avg_views'] = cat_stats['total_clicks'] / cat_stats['vacancy_count']
        cat_stats['avg_applies'] = cat_stats['total_applies'] / cat_stats['vacancy_count']
        cat_stats['conversion_rate'] = (cat_stats['total_applies'] / cat_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
        cat_stats = cat_stats.sort_values('total_clicks', ascending=False)

        media_chart_col, media_commentary_col = st.columns([3, 2])

        with media_chart_col:
            fig_media = go.Figure()
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_views'],
                name='Avg. Views', orientation='h', marker_color=JGP_COLORS['primary']
            ))
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_applies'],
                name='Avg. Applies', orientation='h', marker_color=JGP_COLORS['accent'],
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_media.update_layout(
                barmode='group', height=max(350, len(cat_stats) * 40),
                title="Media Performance by Channel",
                xaxis_title="Average per Vacancy", yaxis_title="",
                legend=dict(orientation='h', y=-0.15),
                plot_bgcolor='rgba(0,0,0,0)',
            )
            st.plotly_chart(fig_media, use_container_width=True)
            report_figures['media'] = fig_media

            # Summary table below chart
            display_media = cat_stats[['source_category', 'vacancy_count', 'avg_views', 'avg_applies', 'conversion_rate']].copy()
            display_media.columns = ['Channel', 'Vacancies', 'Avg. Views', 'Avg. Applies', 'Conversion %']
            display_media['Avg. Views'] = display_media['Avg. Views'].round(1)
            display_media['Avg. Applies'] = display_media['Avg. Applies'].round(1)
            display_media['Conversion %'] = display_media['Conversion %'].round(1)
            st.dataframe(display_media, use_container_width=True, hide_index=True)

        with media_commentary_col:
            media_commentary = generate_section_commentary('media', {
                'cat_stats': cat_stats,
                'client_name': selected_client,
            })
            st.markdown(media_commentary)

        # Source-level detail (expandable)
        with st.expander("View by individual source"):
            media_stats = client_media.groupby(['source_category', 'source']).agg(
                total_clicks=('clicks', 'sum'),
                total_applies=('applies', 'sum'),
                vacancy_count=('entity_id_str', 'nunique')
            ).reset_index()
            media_stats['avg_views'] = media_stats['total_clicks'] / media_stats['vacancy_count']
            media_stats['avg_applies'] = media_stats['total_applies'] / media_stats['vacancy_count']
            media_stats['conversion_rate'] = (media_stats['total_applies'] / media_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
            media_stats = media_stats.sort_values('total_clicks', ascending=False)
            detail_media = media_stats[['source_category', 'source', 'vacancy_count', 'avg_views', 'avg_applies', 'conversion_rate']].copy()
            detail_media.columns = ['Channel', 'Source', 'Vacancies', 'Avg. Views', 'Avg. Applies', 'Conversion %']
            detail_media['Avg. Views'] = detail_media['Avg. Views'].round(1)
            detail_media['Avg. Applies'] = detail_media['Avg. Applies'].round(1)
            detail_media['Conversion %'] = detail_media['Conversion %'].round(1)
            st.dataframe(detail_media, use_container_width=True, hide_index=True)
    else:
        st.info("Media source data not available. Run the `dashboard_media_summary` BigQuery table creation to enable this section.")

    st.markdown("---")

    # ===================================================================
    # POWERPOINT EXPORT
    # ===================================================================
    st.subheader("Export Report")

    # --- Compute additional stats needed for PPTX template ---

    # Slide 2: Top-right quadrant % (vacancies above benchmark on BOTH views and applies)
    benchmarkable_df = scatter_df[scatter_df['category'] == 'Benchmarkable']
    if len(benchmarkable_df) > 0:
        top_quadrant_count = len(benchmarkable_df[
            (benchmarkable_df['views_diff_pct'] > 0) & (benchmarkable_df['applies_diff_pct'] > 0)
        ])
        top_quadrant_pct = (top_quadrant_count / len(benchmarkable_df)) * 100
    else:
        top_quadrant_pct = 0

    # Slide 2: Strongest job category (highest combined diff score, benchmarkable only)
    if len(benchmarkable_df) > 0:
        category_scores = benchmarkable_df.groupby('occupation').agg(
            combined_score=('views_diff_pct', lambda s: s.mean() + benchmarkable_df.loc[s.index, 'applies_diff_pct'].mean()),
            count=('views_diff_pct', 'count')
        ).reset_index()
        # Need at least 2 vacancies in occupation for the category to be considered "strongest"
        category_scores = category_scores[category_scores['count'] >= 2]
        if len(category_scores) > 0:
            top_category = category_scores.sort_values('combined_score', ascending=False).iloc[0]['occupation']
        else:
            top_category = benchmarkable_df.iloc[0]['occupation'] if len(benchmarkable_df) > 0 else 'N/A'
    else:
        top_category = 'N/A'

    # Slide 5: Build the new charts (only if spend entered)
    rate_card_total_val = rate_card_price * num_jobs if annual_spend > 0 else 0
    saving_pct_val = ((rate_card_total_val - annual_spend) / rate_card_total_val * 100) if rate_card_total_val > 0 else 0

    if annual_spend > 0:
        # Stacked bar: Your Spend (bottom) + Saving (top) = Rate Card Total
        saving_amount = max(rate_card_total_val - annual_spend, 0)
        fig_spend_stack = go.Figure()
        fig_spend_stack.add_trace(go.Bar(
            x=['Total Value'],
            y=[annual_spend],
            name='Your Spend',
            marker_color=JGP_COLORS['primary'],
            text=[f"£{annual_spend:,.0f}"],
            textposition='inside',
            textfont=dict(color='white', size=14),
            width=0.7,
        ))
        fig_spend_stack.add_trace(go.Bar(
            x=['Total Value'],
            y=[saving_amount],
            name='Saving vs Rate Card',
            marker_color=JGP_COLORS['accent'],
            text=[f"£{saving_amount:,.0f}"],
            textposition='inside',
            textfont=dict(color=JGP_COLORS['deep_blue'], size=14),
            width=0.7,
        ))
        fig_spend_stack.update_layout(
            barmode='stack',
            title=f"Your Spend vs Rate Card Value (Saving: {saving_pct_val:.0f}%)",
            yaxis_title="GBP",
            height=400,
            bargap=0.05,
            showlegend=True,
            plot_bgcolor='rgba(0,0,0,0)',
            legend=dict(orientation='h', y=-0.15),
        )
        report_figures['spend_vs_ratecard'] = fig_spend_stack

    # Cost per apply by occupation chart (always built when spend > 0)
    cpa_by_occ_fig = None
    if annual_spend > 0:
        roi_by_type_full = client_df.groupby('occupation').agg(
            total_applies=('applies', 'sum'),
            job_count=('clicks', 'count')
        ).reset_index()
        roi_by_type_full = roi_by_type_full[roi_by_type_full['total_applies'] > 0]
        if len(roi_by_type_full) > 0:
            roi_by_type_full['cost_allocated'] = annual_spend * (roi_by_type_full['job_count'] / roi_by_type_full['job_count'].sum())
            roi_by_type_full['cost_per_apply'] = roi_by_type_full['cost_allocated'] / roi_by_type_full['total_applies']
            roi_by_type_full = roi_by_type_full.sort_values('cost_per_apply', ascending=True)

            cpa_by_occ_fig = go.Figure()
            cpa_by_occ_fig.add_trace(go.Bar(
                y=roi_by_type_full['occupation'],
                x=roi_by_type_full['cost_per_apply'],
                orientation='h',
                marker_color=JGP_COLORS['amber'],
                text=roi_by_type_full['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            cpa_by_occ_fig.update_layout(
                title="Cost per Apply by Occupation",
                height=max(350, len(roi_by_type_full) * 32),
                xaxis_title="Cost per Apply (GBP)",
                yaxis_title="",
                plot_bgcolor='rgba(0,0,0,0)',
            )
            report_figures['cost_per_app_by_occupation'] = cpa_by_occ_fig

    # --- Build structured commentary for PPTX template ---
    scatter_struct = generate_section_commentary_structured('benchmark_scatter', {
        'total_count': len(scatter_df),
        'benchmarkable_count': len(scatter_df[scatter_df['category'] == 'Benchmarkable']),
        'zero_applies_count': len(scatter_df[scatter_df['category'].str.startswith('Zero')]),
        'no_benchmark_count': len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)']),
        'top_performers': top_performers,
        'client_name': selected_client,
    })
    average_struct = generate_section_commentary_structured('benchmark_average', {
        'client_avg_clicks': client_avg_clicks,
        'benchmark_avg_clicks': benchmark_avg_clicks,
        'client_avg_applies': client_avg_applies,
        'benchmark_avg_applies': benchmark_avg_applies,
        'num_jobs': len(client_df),
        'client_name': selected_client,
    })
    postings_struct = generate_section_commentary_structured('postings', {
        'total_jobs': len(client_df),
        'total_applies': int(client_df['applies'].sum()),
        'by_type': by_type,
        'client_name': selected_client,
    })
    roi_struct = generate_section_commentary_structured('roi', {
        'annual_spend': annual_spend,
        'rate_card_total': rate_card_total_val,
        'num_jobs': num_jobs,
        'cost_per_apply': annual_spend / total_applies_val if total_applies_val > 0 else 0,
        'saving_pct': saving_pct_val,
        'roi_by_type': roi_by_type_full if (annual_spend > 0 and 'roi_by_type_full' in dir() and roi_by_type_full is not None and len(roi_by_type_full) > 0) else None,
        'client_name': selected_client,
    })
    media_struct = generate_section_commentary_structured('media', {
        'cat_stats': cat_stats,
        'client_name': selected_client,
    })

    # --- Build report_metrics dict (matches template tag names) ---
    report_metrics = {
        # Slide 1
        'client_name': selected_client,
        'PERIOD_START': str(report_start),
        'PERIOD_END': str(report_end),

        # Slide 2 stats
        'stat_total_jobs': f"{num_jobs:,}",
        'stat_top_quadrant_pct': f"{top_quadrant_pct:.0f}",
        'stat_top_category': top_category,

        # Slide 2 commentary
        'commentary_benchmark_intro': scatter_struct['intro'],
        'commentary_benchmark_point_1': scatter_struct['point_1'],
        'commentary_benchmark_point_2': scatter_struct['point_2'],
        'commentary_benchmark_point_3': scatter_struct['point_3'],

        # Slide 3 stats
        'stat_benchmark_average_views': f"{benchmark_avg_clicks:,.0f}",
        'stat_your_jobs_average_views': f"{client_avg_clicks:,.0f}",
        'stat_benchmark_average_applies': f"{benchmark_avg_applies:,.1f}",
        'stat_your_jobs_average_applies': f"{client_avg_applies:,.1f}",

        # Slide 3 commentary
        'commentary_average_intro': average_struct['intro'],
        'commentary_average_point_1': average_struct['point_1'],
        'commentary_average_point_2': average_struct['point_2'],

        # Slide 4 commentary
        'commentary_postings_intro': postings_struct['intro'],
        'commentary_postings_point_1': postings_struct['point_1'],
        'commentary_postings_point_2': postings_struct['point_2'],

        # Slide 5 stats (ROI)
        'stat_cost_per_job': f"£{annual_spend / num_jobs:,.2f}" if (annual_spend > 0 and num_jobs > 0) else "—",
        'stat_cost_per_view': f"£{annual_spend / total_clicks:,.2f}" if (annual_spend > 0 and total_clicks > 0) else "—",
        'stat_cost_per_apply': f"£{annual_spend / total_applies_val:,.2f}" if (annual_spend > 0 and total_applies_val > 0) else "—",

        # Slide 5 commentary
        'commentary_roi_intro': roi_struct['intro'],
        'commentary_roi_point_1': roi_struct['point_1'],
        'commentary_roi_point_2': roi_struct['point_2'],

        # Slide 6 commentary
        'commentary_media_intro': media_struct['intro'],
        'commentary_media_point_1': media_struct['point_1'],
        'commentary_media_point_2': media_struct['point_2'],
        'commentary_media_point_3': media_struct['point_3'],

        # Slide 7 contact
        'contact_name': contact_name or 'Your Account Manager',
        'contact_title': contact_title or 'Account Director',
        'contact_email': contact_email or 'team@jobsgopublic.com',
        'contact_phone': contact_phone or '020 7427 8250',
    }

    # --- Map template chart tags to figures ---
    pptx_figures = {
        'benchmark_scatter': report_figures.get('scatter'),
        'benchmark_average': report_figures.get('benchmark_combined'),
        'postings_by_type': report_figures.get('postings'),
        'spend_vs_ratecard': report_figures.get('spend_vs_ratecard'),
        'cost_per_app_by_occupation': report_figures.get('cost_per_app_by_occupation'),
        'media_performance': report_figures.get('media'),
    }

    # --- Generate PPTX and offer download ---
    template_path = 'Renewals.pptx'
    try:
        pptx_bytes = generate_client_report_pptx(report_metrics, pptx_figures, template_path)
        st.download_button(
            "Download PowerPoint Report",
            data=pptx_bytes,
            file_name=f"advertising_report_{selected_client.replace(' ', '_')}_{report_start}_{report_end}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary"
        )
        st.caption("Download the PowerPoint, edit any commentary you'd like, then File → Export → PDF for the final client copy.")
    except FileNotFoundError:
        st.error(f"Template file not found at `{template_path}`. Make sure `Renewals.pptx` is in the project root.")
    except Exception as e:
        st.warning("PowerPoint generation requires `python-pptx` and `kaleido`.")
        st.caption(f"Error: {type(e).__name__}: {e}")


def generate_client_report_pptx(metrics, figures, template_path):
    """Generate a PowerPoint report by filling a branded template.

    Args:
        metrics: dict mapping placeholder names (without {{}}) to string values
        figures: dict mapping chart slot names (e.g. 'benchmark_scatter') to plotly figures
        template_path: path to the .pptx template file

    Returns: bytes of the populated .pptx file.

    Replaces text placeholders like {{tag_name}} with metrics[tag_name].
    Replaces chart placeholders like {{chart:slot_name}} with PNG images of figures[slot_name].
    """
    from pptx import Presentation
    from pptx.util import Emu
    import io
    import re
    import copy

    prs = Presentation(template_path)

    # --- Helper: render a Plotly figure to PNG bytes at a target aspect ratio ---
    def _fig_to_png(fig, slot_width_emu=None, slot_height_emu=None):
        """Render a Plotly figure to white-background PNG bytes.

        If slot dimensions are passed (the placeholder's .width/.height in
        EMU), the PNG is rendered at that exact aspect ratio so PowerPoint
        doesn't squash the bitmap when fitting it. Otherwise default to 16:9.

        Pixel target is 1500 on the longest side, then kaleido renders at
        scale=4 — effectively very high DPI, sharp at any zoom and when
        exported to PDF. Y/X axes use automargin so tick labels never clip.
        """
        if fig is None:
            return None
        try:
            import plotly.graph_objects as go_local

            # Match the placeholder aspect ratio but keep canvas dimensions
            # constant across charts. Plotly font sizes are pixels relative to
            # canvas, so a uniform canvas keeps text/legend/axes the same
            # relative size on every slide. Crispness comes from `scale` (a
            # pure DPI multiplier), not from a bigger canvas.
            TARGET_LONG_PX = 1800
            if slot_width_emu and slot_height_emu:
                ratio = slot_width_emu / slot_height_emu
            else:
                ratio = 16 / 9
            if ratio >= 1:
                width = TARGET_LONG_PX
                height = int(round(TARGET_LONG_PX / ratio))
            else:
                height = TARGET_LONG_PX
                width = int(round(TARGET_LONG_PX * ratio))

            fig_export = go_local.Figure(fig.to_dict())
            fig_export.update_layout(
                paper_bgcolor='white',
                plot_bgcolor='white',
                font=dict(size=18),
                # pad=15 puts visible space between tick labels and the plot.
                margin=dict(l=120, r=40, t=60, b=80, pad=15),
                # Override the JGP template's small (12pt) legend font for export.
                legend=dict(font=dict(size=18)),
            )
            fig_export.update_xaxes(automargin=True, title_font=dict(size=22))
            fig_export.update_yaxes(automargin=True, title_font=dict(size=22))
            # Force-set bar value-label size and disable auto-shrink. Plotly
            # defaults constraintext='both' which silently shrinks text to fit
            # inside the bar — that's why update_traces alone wasn't visibly
            # bumping the dense postings chart and the narrow spend column.
            for trace in fig_export.data:
                if trace.type != 'bar':
                    continue
                preserved_color = None
                if trace.textfont is not None and trace.textfont.color is not None:
                    preserved_color = trace.textfont.color
                trace.textfont = ({'size': 22, 'color': preserved_color}
                                  if preserved_color else {'size': 22})
                trace.constraintext = 'none'

            # scale=6 → 10800 px on the long side. Maximum crispness; render
            # time is a few seconds per chart, accepted for renewals reports.
            return fig_export.to_image(format='png', width=width, height=height, scale=6)
        except Exception:
            return None

    # --- Helper: replace text in a single shape's text frame, preserving formatting ---
    def _replace_text_in_shape(shape, replacements):
        """Walk runs/paragraphs and replace {{tag}} occurrences. Handles tags split across runs."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            # First try simple per-run replacement (works when tag is fully in one run)
            for run in paragraph.runs:
                text = run.text
                if '{{' in text:
                    for tag, val in replacements.items():
                        placeholder = '{{' + tag + '}}'
                        if placeholder in text:
                            text = text.replace(placeholder, str(val))
                    run.text = text

            # Fallback: if a tag spans multiple runs, the above won't catch it.
            # Concatenate paragraph text, replace, and put it all in the first run.
            full_text = ''.join(r.text for r in paragraph.runs)
            if '{{' in full_text and any('{{' + tag + '}}' in full_text for tag in replacements):
                new_text = full_text
                for tag, val in replacements.items():
                    new_text = new_text.replace('{{' + tag + '}}', str(val))
                if new_text != full_text:
                    if paragraph.runs:
                        paragraph.runs[0].text = new_text
                        for r in paragraph.runs[1:]:
                            r.text = ''

    # --- Step 1: Find all chart placeholders, capture position/size, queue for replacement ---
    # We do this before text replacement so we can find {{chart:xxx}} markers.
    chart_replacements = []  # list of (slide, shape, slot_name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                m = re.search(r'\{\{chart:([^}]+)\}\}', text)
                if m:
                    chart_replacements.append((slide, shape, m.group(1)))

    # --- Step 2: Replace text placeholders on every shape across all slides ---
    # Build the text replacements dict (skip chart tags, those are handled separately)
    text_replacements = {k: v for k, v in metrics.items()}
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_text_in_shape(shape, text_replacements)

    # --- Step 3: Replace chart placeholders with images ---
    for slide, shape, slot_name in chart_replacements:
        fig = figures.get(slot_name)
        # Capture original geometry before deleting the shape
        left, top, width, height = shape.left, shape.top, shape.width, shape.height

        # Remove the placeholder shape
        sp = shape._element
        sp.getparent().remove(sp)

        if fig is not None:
            png_bytes = _fig_to_png(fig, slot_width_emu=width, slot_height_emu=height)
            if png_bytes:
                slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width, height=height)
            else:
                # Failed to render — leave a small note in place
                txt_box = slide.shapes.add_textbox(left, top, width, height)
                txt_box.text_frame.text = '[Chart unavailable]'
        # If fig is None, just remove the placeholder silently

    # --- Step 4: Clean up any remaining unreplaced {{tags}} (set to empty so they don't show) ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join(r.text for r in paragraph.runs)
                    if re.search(r'\{\{[^}]+\}\}', full_text):
                        cleaned = re.sub(r'\{\{[^}]+\}\}', '', full_text)
                        if paragraph.runs:
                            paragraph.runs[0].text = cleaned
                            for r in paragraph.runs[1:]:
                                r.text = ''

    # --- Step 5: Output to bytes ---
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================================
# MAIN APPLICATION
# ============================================================================

def main():
    st.title("📊 Job Performance Dashboard (v2 Dev)")

    # Sidebar - Data Loading Controls
    st.sidebar.header("⚙️ Data Loading Settings")

    days_back = st.sidebar.slider(
        "Vacancies Active In Last (Days)",
        min_value=7,
        max_value=365,
        value=365,
        step=7,
        help="Show vacancies that received at least one view or click within this many days"
    )

    # Sampling option for faster testing
    enable_sampling = st.sidebar.checkbox(
        "Enable Sampling (Faster)",
        value=False,
        help="Limit data to a sample for faster loading during testing"
    )

    sample_size = None
    if enable_sampling:
        sample_size = st.sidebar.number_input(
            "Sample Size (rows)",
            min_value=1000,
            max_value=100000,
            value=10000,
            step=1000,
            help="Number of rows to sample from BigQuery"
        )

    # Load data with progress indicators
    progress_bar = st.progress(0)
    status_text = st.empty()

    status_text.text("Loading data from BigQuery... 0%")
    df_raw, daily_totals, media_df, region_raw = load_all_data(days_back=days_back, sample_size=sample_size)
    progress_bar.progress(30)

    status_text.text("Loading launch timing data... 30%")
    launch_timing_df = load_launch_timing_data(days_back=days_back)
    progress_bar.progress(40)

    status_text.text("Loading importer mapping... 40%")
    importer_mapping = load_importer_mapping()
    progress_bar.progress(50)

    # Process enriched data
    status_text.text("Preparing enriched data... 50%")
    df = df_raw.copy()
    df = prepare_enriched_data(df)  # Rename enriched table columns
    progress_bar.progress(60)

    status_text.text("Applying importer mapping... 60%")
    df = apply_importer_mapping(df, importer_mapping)
    progress_bar.progress(70)

    status_text.text("Parsing upgrades... 70%")
    df = parse_upgrades(df)
    progress_bar.progress(80)

    status_text.text("Parsing dates... 80%")
    df = parse_dates_in_jobiqo(df)  # Parse timestamp columns
    progress_bar.progress(90)

    status_text.text("Adding occupation column... 88%")
    df = add_occupation_column(df)
    progress_bar.progress(90)

    status_text.text("Processing salary data... 92%")
    df = process_salary_columns(df)
    progress_bar.progress(100)

    # Process region_df through same enrichment pipeline
    region_df = None
    if region_raw is not None:
        region_df = region_raw.copy()
        region_df = prepare_enriched_data(region_df)
        region_df = apply_importer_mapping(region_df, importer_mapping)
        region_df = parse_upgrades(region_df)
        region_df = parse_dates_in_jobiqo(region_df)
        region_df = add_occupation_column(region_df)
        region_df = process_salary_columns(region_df)
    st.session_state['_region_df'] = region_df

    status_text.text("✅ Data loaded successfully!")
    progress_bar.empty()
    status_text.empty()

    # Initialize session state for all tabs
    for tab_prefix in ['overview', 'deepdive', 'vacancy', 'comp_left', 'comp_right', 'sales', 'salary']:
        if f'{tab_prefix}_filters' not in st.session_state:
            st.session_state[f'{tab_prefix}_filters'] = None

    # Sidebar
    st.sidebar.markdown("---")
    st.sidebar.header("📊 Dashboard Info")

    # Authentication status
    try:
        if hasattr(st, 'secrets') and 'gcp_service_account' in st.secrets:
            st.sidebar.success("🔐 Authentication: Streamlit Secrets")
        else:
            st.sidebar.info("🔐 Authentication: Local File")
    except:
        st.sidebar.info("🔐 Authentication: Local File")

    # Data loading info
    st.sidebar.info(f"📊 Vacancies active in last {days_back} days")
    st.sidebar.metric("Total Vacancies", f"{len(df):,}")
    if 'clicks' in df.columns:
        st.sidebar.metric("Total Clicks", f"{int(df['clicks'].sum()):,}")
        st.sidebar.metric("Total Applies", f"{int(df['applies'].sum()):,}")
    st.sidebar.info(f"Last updated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    if enable_sampling:
        st.sidebar.warning(f"⚠️ Sampling enabled: showing {len(df):,} of all vacancies")

    if st.sidebar.button("🔄 Refresh Data"):
        st.cache_data.clear()
        st.rerun()

    # Debug: Show importer mapping status
    with st.sidebar.expander("🔧 Importer Mapping Debug"):
        st.write(f"Mapping loaded: {len(importer_mapping)} entries")
        if importer_mapping:
            st.write("**Mapping:**")
            for k, v in importer_mapping.items():
                st.write(f"'{k}' → '{v}'")

        # Show unique importer IDs in data
        if 'importer_ID' in df.columns:
            unique_ids = df['importer_ID'].dropna().astype(str).str.strip().unique()
            st.write(f"\n**Unique IDs in data:** {len(unique_ids)}")
            for uid in sorted(unique_ids, key=str):
                matched = importer_mapping.get(uid, "NOT FOUND")
                st.write(f"'{uid}' → {matched}")

        # Show unique importer names in data
        if 'importer_name' in df.columns:
            unique_names = df['importer_name'].dropna().astype(str).unique()
            st.write(f"\n**Unique names in data:** {len(unique_names)}")
            for name in sorted(unique_names, key=str):
                count = len(df[df['importer_name'] == name])
                st.write(f"'{name}': {count} vacancies")

    # Tabs
    tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8 = st.tabs([
        "📊 Overview",
        "🔍 Deep Dive",
        "📋 Vacancy Performance",
        "⚖️ Comparison",
        "💼 Sales Intelligence",
        "📅 Launch Timing",
        "📄 Client Report",
        "💰 Salary Benchmarking"
    ])

    with tab1:
        create_overview_tab(df, daily_totals=daily_totals)

    with tab2:
        create_deep_dive_tab(df)

    with tab3:
        create_vacancy_performance_tab(df, full_df=df)

    with tab4:
        create_comparison_tab(df)

    with tab5:
        create_sales_intelligence_tab(df)

    with tab6:
        create_launch_timing_tab(launch_timing_df)

    with tab7:
        create_client_report_tab(df, media_df=media_df)

    with tab8:
        render_salary(df)

if __name__ == "__main__":
    main()
